"""Generate VIGA Project Summary PowerPoint from docs/ folder.

Layout per document type:
  Run:      Page 1 = overview (2 large images: input+output)
            Page 2+ = round outputs (6 small images per page)
  Analysis: 1+ pages: text page + large image pages (2 per page)
  Summary:  2+ pages: text page(s) + large image pages (2 per page)
"""
import io
import os
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DOCS = Path("docs")
OUT_PART1 = Path("docs/VIGA_Project_Summary_v5_part1.pptx")
OUT_PART2 = Path("docs/VIGA_Project_Summary_v5_part2.pptx")
OUT_SAM3D = Path("docs/VIGA_SAM3D_Feb15-24.pptx")

# Max pixels per inch for embedded images (used only in compressed mode)
EMBED_DPI = 150

# Global flag: when True, non-GIF images are downscaled before embedding
COMPRESS = False

# Colors (dark theme)
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
BG_CARD = RGBColor(0x16, 0x1B, 0x22)
TEXT_PRIMARY = RGBColor(0xE6, 0xED, 0xF3)
TEXT_SECONDARY = RGBColor(0x8B, 0x94, 0x9E)
ACCENT_BLUE = RGBColor(0x58, 0xA6, 0xFF)
ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_PURPLE = RGBColor(0xA3, 0x71, 0xF7)
ACCENT_ORANGE = RGBColor(0xD2, 0x99, 0x22)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

GITHUB_DOCS_BASE = "https://github.com/frankwings/VIGA_exploration/blob/main/docs/"


# ============================================================================
# Helpers
# ============================================================================

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18,
             color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    tb.text_frame.vertical_anchor = vertical_anchor
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = align
    return tb


def add_multiline(slide, left, top, width, height, lines, size=12, color=TEXT_PRIMARY):
    """Add multiple lines of text (each line = a paragraph)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(3)
    return tb


def add_hyperlink(slide, left, top, width, height, text, url,
                  size=10, color=ACCENT_BLUE):
    """Add a textbox with a clickable hyperlink."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = "Segoe UI"
    r.font.underline = True
    r.hyperlink.address = url
    return tb


def add_rect(slide, left, top, width, height, fill_color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s


def _shrink_png(full_path, display_w_inches, display_h_inches):
    """Downscale a non-GIF image to fit display size at EMBED_DPI. Returns BytesIO PNG."""
    max_px_w = int(display_w_inches * EMBED_DPI) if display_w_inches else 2000
    max_px_h = int(display_h_inches * EMBED_DPI) if display_h_inches else 2000
    with Image.open(str(full_path)) as img:
        img = img.convert("RGBA") if img.mode == "RGBA" else img.convert("RGB")
        native_w, native_h = img.size
        if native_w > max_px_w or native_h > max_px_h:
            img.thumbnail((max_px_w, max_px_h), Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="PNG", optimize=True)
        buf.seek(0)
    return buf, native_w, native_h


def add_img(slide, img_path, left, top, width=None, height=None):
    """Add image preserving aspect ratio, fitted within width x height box.

    When COMPRESS is True, non-GIF images are downscaled to EMBED_DPI.
    GIFs are always embedded as-is (animated) regardless of COMPRESS setting.
    """
    full = DOCS / img_path
    if not full.exists():
        return False
    try:
        is_gif = full.suffix.lower() == ".gif"

        if COMPRESS and not is_gif:
            disp_w_in = width / 914400 if width else None
            disp_h_in = height / 914400 if height else None
            source, native_w, native_h = _shrink_png(full, disp_w_in, disp_h_in)
        else:
            # Use original file — read native dimensions for aspect ratio
            with Image.open(str(full)) as img:
                native_w, native_h = img.size
            source = str(full)

        if width and height:
            aspect = native_w / native_h
            box_aspect = width / height
            if aspect >= box_aspect:
                final_w = width
                final_h = int(width / aspect)
            else:
                final_h = height
                final_w = int(height * aspect)
            x_off = (width - final_w) // 2
            y_off = (height - final_h) // 2
            slide.shapes.add_picture(source, left + x_off, top + y_off,
                                     final_w, final_h)
        elif width:
            slide.shapes.add_picture(source, left, top, width=width)
        elif height:
            slide.shapes.add_picture(source, left, top, height=height)
        else:
            slide.shapes.add_picture(source, left, top)
        return True
    except Exception:
        return False


# ============================================================================
# Page builders
# ============================================================================

def make_title_slide(prs, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    add_text(s, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "VIGA", 60, ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.0), Inches(11), Inches(1),
             "Vision-as-Inverse-Graphics Agent", 28, TEXT_PRIMARY, False, PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
             "Iterative Generate / Render / Verify Loop using VLMs + Blender",
             18, TEXT_SECONDARY, False, PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             subtitle or "Project Summary  |  January - February 2026",
             16, TEXT_SECONDARY, False, PP_ALIGN.CENTER)


def make_flow_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Pipeline Architecture", 36, ACCENT_BLUE, True)

    # ---- Left Column: Asset Sources ----
    add_text(s, Inches(0.4), Inches(0.85), Inches(4.0), Inches(0.3),
             "Asset Sources", 16, ACCENT_GREEN, True)

    # Target Image
    add_rect(s, Inches(0.4), Inches(1.15), Inches(4.0), Inches(0.5), BG_CARD)
    add_text(s, Inches(0.4), Inches(1.2), Inches(4.0), Inches(0.4),
             "Target Image  (photograph / artwork)", 12, ACCENT_ORANGE, True,
             PP_ALIGN.CENTER)

    # Arrow
    add_text(s, Inches(0.4), Inches(1.6), Inches(4.0), Inches(0.2),
             "\u25bc", 12, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    # SAM3D
    add_rect(s, Inches(0.4), Inches(1.8), Inches(4.0), Inches(1.3), BG_CARD)
    add_text(s, Inches(0.6), Inches(1.85), Inches(3.6), Inches(0.3),
             "SAM3D  (image \u2192 3D)", 13, ACCENT_PURPLE, True)
    add_multiline(s, Inches(0.6), Inches(2.15), Inches(3.6), Inches(0.9), [
        "1. SAM ViT-H \u2192 per-object binary masks",
        "2. TRELLIS \u2192 3D mesh per masked region",
        "Out: *.glb (vertex colors) + transform JSON",
        "~2 hours for 8 objects (RTX 5080 GPU)",
    ], 10, TEXT_SECONDARY)

    # Arrow
    add_text(s, Inches(0.4), Inches(3.05), Inches(4.0), Inches(0.2),
             "\u25bc  GLBs + transforms \u2192 Generator", 9, TEXT_SECONDARY,
             False, PP_ALIGN.CENTER)

    # Meshy
    add_rect(s, Inches(0.4), Inches(3.25), Inches(4.0), Inches(1.3), BG_CARD)
    add_text(s, Inches(0.6), Inches(3.3), Inches(3.6), Inches(0.3),
             "Meshy  (text \u2192 3D)", 13, ACCENT_ORANGE, True)
    add_multiline(s, Inches(0.6), Inches(3.6), Inches(3.6), Inches(0.9), [
        "1. Fuzzy-match local cache (instant)",
        "2. API fallback: preview \u2192 refine \u2192 GLB",
        "Out: *.glb (UV textured, high quality)",
        "~0s cached  /  ~7 min per object via API",
    ], 10, TEXT_SECONDARY)

    # ---- Right Column: Core Loop ----
    add_text(s, Inches(4.7), Inches(0.85), Inches(8.2), Inches(0.3),
             "Core Loop  (\u226425 rounds)", 16, ACCENT_GREEN, True)

    # Generator Agent
    add_rect(s, Inches(4.7), Inches(1.15), Inches(8.2), Inches(1.35), BG_CARD)
    add_text(s, Inches(4.9), Inches(1.2), Inches(5.0), Inches(0.3),
             "Generator Agent  (VLM: GPT-5)", 13, ACCENT_GREEN, True)
    add_multiline(s, Inches(4.9), Inches(1.5), Inches(7.8), Inches(0.9), [
        "Writes Blender Python: import GLBs, compose scene, lighting + camera",
        "Selects best GLB per object: SAM3D (shape) vs Meshy (texture)",
        "Tools: initialize_plan | get_better_object | execute_and_evaluate | end",
        "Memory: target image + SAM3D paths + transforms + Verifier feedback",
    ], 10, TEXT_SECONDARY)

    # Arrow
    add_text(s, Inches(4.7), Inches(2.45), Inches(8.2), Inches(0.2),
             "\u25bc  Blender Python script", 9, TEXT_SECONDARY, False,
             PP_ALIGN.CENTER)

    # Executor (left half of right column)
    add_rect(s, Inches(4.7), Inches(2.65), Inches(3.9), Inches(1.35), BG_CARD)
    add_text(s, Inches(4.9), Inches(2.7), Inches(3.5), Inches(0.3),
             "Blender Executor  (EEVEE)", 13, ACCENT_BLUE, True)
    add_multiline(s, Inches(4.9), Inches(3.0), Inches(3.5), Inches(0.9), [
        "Headless: blender --background",
        "Wrapper sets EEVEE Next engine",
        "Runs Generator's Python script",
        "Out: Camera*.png + state.blend",
    ], 10, TEXT_SECONDARY)

    # Arrow between executor and verifier
    add_text(s, Inches(8.55), Inches(3.1), Inches(0.5), Inches(0.4),
             "\u2192", 20, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    # Verifier (right half of right column)
    add_rect(s, Inches(9.0), Inches(2.65), Inches(3.9), Inches(1.35), BG_CARD)
    add_text(s, Inches(9.2), Inches(2.7), Inches(3.5), Inches(0.3),
             "Verifier Agent  (VLM)", 13, ACCENT_BLUE, True)
    add_multiline(s, Inches(9.2), Inches(3.0), Inches(3.5), Inches(0.9), [
        "Compares render \u2194 target image",
        "Checks: layout, materials, scale",
        "Tools: rotate, set_camera, scene_info",
        "Out: structured text feedback",
    ], 10, TEXT_SECONDARY)

    # Feedback loop text
    add_text(s, Inches(4.7), Inches(3.95), Inches(8.2), Inches(0.3),
             "\u21bb  Feedback \u2192 Generator memory \u2192 next round"
             "  (until approved or max rounds)",
             10, ACCENT_PURPLE, False, PP_ALIGN.CENTER)

    # ---- Bottom: Outputs + Modes ----
    # Output Structure
    add_rect(s, Inches(0.4), Inches(4.55), Inches(6.0), Inches(2.7), BG_CARD)
    add_text(s, Inches(0.6), Inches(4.6), Inches(5.6), Inches(0.3),
             "Output Structure", 13, ACCENT_GREEN, True)
    add_multiline(s, Inches(0.6), Inches(4.9), Inches(5.6), Inches(2.2), [
        "output/{mode}/{timestamp}/{task}/",
        "  scripts/{N}.py           Blender Python per round",
        "  renders/{N}/Camera*.png  rendered frames",
        "  renders/{N}/state.blend  scene snapshot (undo)",
        "  generator_memory.json    full LLM conversation",
        "  verifier_memory.json     all feedback rounds",
        "  blender_file.blend       final Blender scene",
        "",
        "data/{mode}/{task}/assets/   Meshy GLB cache",
    ], 10, TEXT_SECONDARY)

    # Pipeline Modes
    add_rect(s, Inches(6.7), Inches(4.55), Inches(6.2), Inches(2.7), BG_CARD)
    add_text(s, Inches(6.9), Inches(4.6), Inches(5.8), Inches(0.3),
             "Pipeline Modes", 13, ACCENT_GREEN, True)
    add_multiline(s, Inches(6.9), Inches(4.9), Inches(5.8), Inches(2.2), [
        "get_asset",
        "  Meshy only \u2014 text-to-3D, no SAM3D",
        "",
        "get_asset_sam3d",
        "  SAM3D auto-init + Meshy replaces bad objects",
        "",
        "Pre-computed (--sam3d-results)",
        "  Load existing SAM3D GLBs, skip reconstruction",
        "",
        "Static: Camera.png  |  Dynamic: Camera_f{NNNN}.png",
    ], 10, TEXT_SECONDARY)


def make_env_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Environment Setup", 36, ACCENT_BLUE, True)
    add_rect(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.8), BG_CARD)
    add_text(s, Inches(0.7), Inches(1.3), Inches(5.6), Inches(0.5),
             "Windows 11 (Primary)", 20, ACCENT_GREEN, True)
    add_multiline(s, Inches(0.7), Inches(1.9), Inches(5.6), Inches(5.0), [
        "OS: Windows 11 Home 10.0.26200",
        "GPU: NVIDIA RTX 5080 16GB VRAM",
        "CPU: AMD Ryzen 9 9900X | 32GB DDR5-6000",
        "",
        "Conda Env: agent (Python 3.10)",
        "  openai 2.6.1  |  mcp 1.20.0",
        "  pillow 12.0  |  numpy 2.2.6",
        "  python-pptx 1.0.2",
        "",
        "Conda Env: sam (Python 3.10)",
        "  torch 2.10.0+cu128 (CUDA 12.8)",
        "  torchvision 0.25.0+cu128",
        "  segment-anything 1.0 (Meta SAM)",
        "  opencv-python 4.13.0",
        "",
        "Blender 4.5",
        "  EEVEE (BLENDER_EEVEE_NEXT)",
        "  Cycles GPU rendering",
    ], 11, TEXT_PRIMARY)
    add_rect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.8), BG_CARD)
    add_text(s, Inches(7.0), Inches(1.3), Inches(5.6), Inches(0.5),
             "Linux / WSL2 (Tested)", 20, ACCENT_GREEN, True)
    add_multiline(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(5.0), [
        "WSL2 Ubuntu (SAM3D testing)",
        "  CUDA passthrough from host GPU",
        "  ~30% faster SAM3D inference",
        "  Higher VRAM efficiency",
        "",
        "Windows Compatibility Fixes:",
        "  Path resolution: .resolve() all Paths",
        "  Pipe deadlock: temp files (not pipes)",
        "  Encoding: always utf-8 (not cp1252)",
        "  Spaces in paths: quote subprocess args",
        "",
        "VLM Model: GPT-5 (OpenAI API)",
        "  Generator: writes Blender Python",
        "  Verifier: evaluates renders vs target",
        "",
        "3D Asset Sources:",
        "  SAM3D: image -> 3D mesh (TRELLIS)",
        "  Meshy: text -> 3D mesh (API + cache)",
    ], 11, TEXT_PRIMARY)


def make_text_page(prs, date_str, author, badge, badge_color, title, summary, points):
    """A page with header, summary text, and bullet points."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    add_text(s, Inches(0.5), Inches(0.2), Inches(9.5), Inches(0.7),
             f"{date_str}  \u2014  {title}", 28, ACCENT_BLUE, True)
    add_text(s, Inches(10.2), Inches(0.3), Inches(2.8), Inches(0.5),
             f"{author}  |  {badge}", 14, badge_color, False, PP_ALIGN.RIGHT)
    y = Inches(1.0)
    if summary:
        add_text(s, Inches(0.5), y, Inches(12.3), Inches(1.1),
                 summary, 16, TEXT_SECONDARY)
        y = Inches(2.2)
    if points:
        add_multiline(s, Inches(0.7), y, Inches(11.9), Inches(7.5) - y - Inches(0.2),
                      points, 16, TEXT_PRIMARY)
    return s


def make_large_image_pages(prs, date_str, title, images):
    """Pages with 2 large images each, side by side, maximized."""
    if not images:
        return
    for i in range(0, len(images), 2):
        batch = images[i:i + 2]
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(s, BG_DARK)
        pg = f" ({i // 2 + 1}/{-(-len(images) // 2)})" if len(images) > 2 else ""
        add_text(s, Inches(0.2), Inches(0.08), Inches(12.9), Inches(0.4),
                 f"{date_str}  \u2014  {title}{pg}", 20, ACCENT_BLUE, True)
        if len(batch) == 1:
            # Single image — center it large
            label, path = batch[0]
            w = Inches(10.0)
            add_text(s, Inches(1.6), Inches(0.5), w, Inches(0.25),
                     label, 11, TEXT_SECONDARY, True, PP_ALIGN.CENTER)
            add_img(s, path, Inches(1.6), Inches(0.8), width=w, height=Inches(6.4))
        else:
            # Two images side by side — maximized
            w = Inches(6.4)
            positions = [Inches(0.15), Inches(6.75)]
            for j, (label, path) in enumerate(batch):
                x = positions[j]
                add_text(s, x, Inches(0.5), w, Inches(0.25),
                         label, 11, TEXT_SECONDARY, True, PP_ALIGN.CENTER)
                add_img(s, path, x, Inches(0.8), width=w, height=Inches(6.4))


def make_run_overview(prs, date_str, author, title, summary, input_img, output_img):
    """Run page 1: overview with input/output images, maximized."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    add_text(s, Inches(0.2), Inches(0.1), Inches(9.8), Inches(0.5),
             f"{date_str}  \u2014  {title}", 24, ACCENT_BLUE, True)
    add_text(s, Inches(10.2), Inches(0.15), Inches(2.9), Inches(0.4),
             f"{author}  |  Run", 11, ACCENT_GREEN, False, PP_ALIGN.RIGHT)
    add_text(s, Inches(0.2), Inches(0.65), Inches(12.9), Inches(0.55),
             summary, 11, TEXT_SECONDARY)
    img_w = Inches(6.2)
    img_h = Inches(5.5)
    y_label = Inches(1.25)
    y_img = Inches(1.55)
    add_text(s, Inches(0.3), y_label, img_w, Inches(0.25),
             "Input / Target", 12, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
    if input_img:
        add_img(s, input_img, Inches(0.3), y_img, width=img_w, height=img_h)
    add_text(s, Inches(6.8), y_label, img_w, Inches(0.25),
             "Output / Result", 12, ACCENT_GREEN, True, PP_ALIGN.CENTER)
    if output_img:
        add_img(s, output_img, Inches(6.8), y_img, width=img_w, height=img_h)


def make_run_rounds(prs, date_str, title, rounds):
    """Run page 2+: round outputs, 6 per page in 3 cols x 2 rows, maximized."""
    if not rounds:
        return
    per_page = 6
    for idx in range(0, len(rounds), per_page):
        batch = rounds[idx:idx + per_page]
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(s, BG_DARK)
        pg = f" (page {idx // per_page + 1})" if len(rounds) > per_page else ""
        add_text(s, Inches(0.2), Inches(0.08), Inches(12.9), Inches(0.38),
                 f"{date_str}  \u2014  {title}  \u2014  Rounds{pg}", 18, ACCENT_BLUE, True)
        # 3 cols x 2 rows — maximized
        col_w = Inches(4.2)
        img_h = Inches(3.1)
        x_pos = [Inches(0.2), Inches(4.55), Inches(8.9)]
        y_rows = [Inches(0.5), Inches(3.9)]
        for i, (label, path) in enumerate(batch):
            col = i % 3
            row = i // 3
            x = x_pos[col]
            y = y_rows[row]
            add_text(s, x, y, col_w, Inches(0.2),
                     label, 9, TEXT_SECONDARY, True, PP_ALIGN.CENTER)
            add_img(s, path, x, y + Inches(0.22), width=col_w, height=img_h)


def make_pipeline_table_slide(prs, date_str, title):
    """SAM3D Pipeline: fancy table showing 6 steps with colors."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)

    # Title
    add_text(s, Inches(0.2), Inches(0.1), Inches(12.9), Inches(0.5),
             f"{date_str}  —  {title}", 24, ACCENT_BLUE, True)

    # Table data
    rows = [
        ("Step", "Model / Env", "Input", "Output", "Key Detail"),
        ("1", "SAM ViT-H (sam)", "Image", "N masks", "Binary (H,W) uint8"),
        ("2", "MoGe (sam3d_py311)", "Image", "Pointmap + intrinsics", "(H,W,3) camera-space 3D coords"),
        ("3", "TRELLIS (sam3d_py311)", "RGBA per mask", "3D mesh", "SS (2) + SLAT (12) + decoder"),
        ("4", "SS Model", "Sparse struct + pointmap", "S, R, T", "Initial pose prediction"),
        ("5", "Layout Optimizer", "GLB + pointmap", "Refined S, R, T", "Silhouette rendering (5a-5c)"),
        ("6", "Export", "Posed mesh", "GLB (PyTorch3D space)", "Baked transforms + vertex colors"),
    ]

    # Dimensions with larger fonts
    left = Inches(0.2)
    top = Inches(0.75)
    col_widths = [Inches(0.6), Inches(2.2), Inches(2.2), Inches(2.2), Inches(4.5)]
    row_height = Inches(0.65)  # Increased from 0.45

    # Draw table
    for row_idx, row in enumerate(rows):
        is_header = (row_idx == 0)
        bg_color = BG_CARD if is_header else BG_DARK
        text_color = ACCENT_GREEN if is_header else TEXT_SECONDARY
        text_size = 14 if is_header else 12  # Increased from 11 and 10

        x = left
        for col_idx, (col_width, cell_text) in enumerate(zip(col_widths, row)):
            # Cell background
            add_rect(s, x, top + row_idx * row_height, col_width, row_height,
                    BG_CARD if is_header else BG_DARK)

            # Cell border (simplified: just right border)
            if col_idx < len(col_widths) - 1:
                line = s.shapes.add_connector(1, x + col_width, top + row_idx * row_height,
                                            x + col_width, top + (row_idx + 1) * row_height)
                line.line.color.rgb = TEXT_SECONDARY
                line.line.width = Pt(0.5)

            # Cell text centered both horizontally and vertically
            add_text(s, x + Inches(0.05), top + row_idx * row_height,
                    col_width - Inches(0.1), row_height,
                    cell_text, text_size, text_color, is_header, PP_ALIGN.CENTER)
            x += col_width


def make_pipeline_flowchart_slide(prs, date_str, title):
    """SAM3D Pipeline: large-font flowchart with single-line text, filling entire slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)

    # Title
    add_text(s, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.55),
             "SAM3D Image-to-3D Pipeline", 40, ACCENT_BLUE, True)

    # Left column
    add_text(s, Inches(0.3), Inches(0.75), Inches(3.0), Inches(0.3),
             "Input & Segmentation", 16, ACCENT_GREEN, True)

    # Input
    add_rect(s, Inches(0.3), Inches(1.1), Inches(3.0), Inches(0.65), BG_CARD)
    add_text(s, Inches(0.4), Inches(1.2), Inches(2.8), Inches(0.45),
             "Input: Scene Image", 18, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

    add_text(s, Inches(0.3), Inches(1.78), Inches(3.0), Inches(0.25),
             "▼", 20, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    # Step 1: SAM
    add_rect(s, Inches(0.3), Inches(2.08), Inches(3.0), Inches(1.0), BG_CARD)
    add_text(s, Inches(0.4), Inches(2.15), Inches(2.8), Inches(0.32),
             "Step 1: SAM Segmentation", 16, ACCENT_PURPLE, True, PP_ALIGN.CENTER)
    add_text(s, Inches(0.4), Inches(2.5), Inches(2.8), Inches(0.45),
             "N binary masks (H, W)", 14, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    add_text(s, Inches(0.3), Inches(3.12), Inches(3.0), Inches(0.25),
             "▼", 20, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    # Step 2: MoGe
    add_rect(s, Inches(0.3), Inches(3.42), Inches(3.0), Inches(1.0), BG_CARD)
    add_text(s, Inches(0.4), Inches(3.49), Inches(2.8), Inches(0.32),
             "Step 2: MoGe Depth", 16, ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_text(s, Inches(0.4), Inches(3.84), Inches(2.8), Inches(0.45),
             "Pointmap (H,W,3) + intrinsics", 14, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    # Middle column
    add_text(s, Inches(3.6), Inches(0.75), Inches(3.2), Inches(0.3),
             "3D Reconstruction", 16, ACCENT_GREEN, True)

    # Step 3: TRELLIS
    add_rect(s, Inches(3.6), Inches(1.1), Inches(3.2), Inches(1.0), BG_CARD)
    add_text(s, Inches(3.7), Inches(1.17), Inches(3.0), Inches(0.32),
             "Step 3: TRELLIS 3D", 16, ACCENT_GREEN, True, PP_ALIGN.CENTER)
    add_text(s, Inches(3.7), Inches(1.52), Inches(3.0), Inches(0.45),
             "80K verts, vertex colors", 14, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    add_text(s, Inches(3.6), Inches(2.15), Inches(3.2), Inches(0.25),
             "▼", 20, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    # Step 4+5: Pose
    add_rect(s, Inches(3.6), Inches(2.45), Inches(3.2), Inches(2.0), BG_CARD)
    add_text(s, Inches(3.7), Inches(2.52), Inches(3.0), Inches(0.32),
             "Step 4: SS Model Pose", 16, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
    add_text(s, Inches(3.7), Inches(2.88), Inches(3.0), Inches(0.35),
             "Predict S, R, T", 14, TEXT_SECONDARY, False, PP_ALIGN.CENTER)
    add_text(s, Inches(3.7), Inches(3.28), Inches(3.0), Inches(0.32),
             "Step 5: Optimize Layout", 16, ACCENT_PURPLE, True, PP_ALIGN.CENTER)
    add_text(s, Inches(3.7), Inches(3.64), Inches(3.0), Inches(0.35),
             "Refine with ICP + rendering", 14, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    add_text(s, Inches(3.6), Inches(4.48), Inches(3.2), Inches(0.25),
             "▼", 20, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    # Step 6: Export
    add_rect(s, Inches(3.6), Inches(4.78), Inches(3.2), Inches(0.85), BG_CARD)
    add_text(s, Inches(3.7), Inches(4.85), Inches(3.0), Inches(0.32),
             "Step 6: Export GLB", 16, ACCENT_GREEN, True, PP_ALIGN.CENTER)
    add_text(s, Inches(3.7), Inches(5.2), Inches(3.0), Inches(0.35),
             "PyTorch3D camera space", 14, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    # Right column: Key insights
    add_text(s, Inches(7.2), Inches(0.75), Inches(5.5), Inches(0.3),
             "Key Technical Details", 16, ACCENT_GREEN, True)

    add_rect(s, Inches(7.2), Inches(1.1), Inches(5.5), Inches(1.8), BG_CARD)
    add_text(s, Inches(7.3), Inches(1.18), Inches(5.3), Inches(0.32),
             "MoGe = Single Source of Camera Calibration", 14, ACCENT_ORANGE, True, PP_ALIGN.LEFT)
    add_text(s, Inches(7.3), Inches(1.53), Inches(5.3), Inches(1.3),
             "Wrong intrinsics cascade to all downstream stages", 12, TEXT_SECONDARY, False, PP_ALIGN.LEFT)

    add_rect(s, Inches(7.2), Inches(3.0), Inches(5.5), Inches(1.8), BG_CARD)
    add_text(s, Inches(7.3), Inches(3.08), Inches(5.3), Inches(0.32),
             "TRELLIS Receives ONLY RGBA Image", 14, ACCENT_PURPLE, True, PP_ALIGN.LEFT)
    add_text(s, Inches(7.3), Inches(3.43), Inches(5.3), Inches(1.3),
             "NOT depth. 3D shape from appearance only. MoGe used for pose anchoring", 12, TEXT_SECONDARY, False, PP_ALIGN.LEFT)

    add_rect(s, Inches(7.2), Inches(4.9), Inches(5.5), Inches(0.73), BG_CARD)
    add_text(s, Inches(7.3), Inches(4.97), Inches(5.3), Inches(0.32),
             "Coordinate Space Flow", 14, ACCENT_BLUE, True, PP_ALIGN.LEFT)
    add_text(s, Inches(7.3), Inches(5.32), Inches(5.3), Inches(0.25),
             "MoGe → PyTorch3D → Blender GLTF", 12, TEXT_SECONDARY, False, PP_ALIGN.LEFT)


def make_architecture_table_slide(prs, date_str, title):
    """TRELLIS v1 vs SAM3D vs TRELLIS.2: fancy comparison table."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)

    # Title
    add_text(s, Inches(0.2), Inches(0.1), Inches(12.9), Inches(0.4),
             f"{date_str}  —  {title}", 22, ACCENT_BLUE, True)

    # Table data with 4 columns
    rows = [
        ("Aspect", "TRELLIS v1 (Base)", "SAM3D (T1 + Meta)", "TRELLIS.2"),
        ("Backbone", "SLAT dense flow", "TRELLIS v1 backbone", "O-Voxel octree"),
        ("Model Size", "~500M params", "~500M (T1) + heads", "~4B params"),
        ("Resolution", "SLAT 64³ dense", "SLAT 64³ dense", "Up to 1536³ octree"),
        ("Architecture", "SS + SLAT + decoder", "SS + SLAT + MoT poses", "3-stage DiT"),
        ("Pose Estimation", "Layout optimizer only", "SS model (end-to-end)", "None built-in"),
        ("Pose Prior", "No SS model", "MoT heads added", "No equivalent"),
        ("VRAM Required", "~3 GB", "~3-4 GB (7 models)", "~9-10 GB"),
        ("Speed (per object)", "~60-90s (L4)", "~60-90s (L4)", "~180s (L4)"),
        ("Mesh Quality", "80K verts, vert colors", "80K verts, vert colors", "500K-2.4M verts, PBR"),
        ("Intended Use", "Baseline 3D recon", "Scene understanding", "High-quality meshes"),
        ("IoU (Dining)", "N/A (baseline)", "0.53 (reference)", "0.23 (w/o SS)"),
    ]

    # Dimensions: fits within slide (13" wide x 7.5" tall)
    left = Inches(0.15)
    top = Inches(0.6)
    col_widths = [Inches(2.1), Inches(3.2), Inches(3.2), Inches(3.2)]  # Total 11.7" + margins
    row_height = Inches(0.54)  # Fits all 12 rows on one slide

    # Draw table
    for row_idx, row in enumerate(rows):
        is_header = (row_idx == 0)
        bg_color = BG_CARD if is_header else BG_DARK
        text_color = ACCENT_GREEN if is_header else TEXT_SECONDARY
        text_size = 16 if is_header else 12  # Increased from 10 and 8.5

        x = left
        for col_idx, (col_width, cell_text) in enumerate(zip(col_widths, row)):
            # Cell background
            add_rect(s, x, top + row_idx * row_height, col_width, row_height,
                    BG_CARD if is_header else BG_DARK)

            # Cell border (simplified: just right border)
            if col_idx < len(col_widths) - 1:
                line = s.shapes.add_connector(1, x + col_width, top + row_idx * row_height,
                                            x + col_width, top + (row_idx + 1) * row_height)
                line.line.color.rgb = TEXT_SECONDARY
                line.line.width = Pt(0.5)

            # Cell text centered both horizontally and vertically
            add_text(s, x + Inches(0.05), top + row_idx * row_height,
                    col_width - Inches(0.1), row_height,
                    cell_text, text_size, text_color, is_header, PP_ALIGN.CENTER)
            x += col_width


def make_closing_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    add_text(s, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
             "VIGA Project Summary", 48, ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.5), Inches(11), Inches(1),
             "52 docs  |  27 days  |  SAM3D + Meshy + Blender + GPT-5",
             20, TEXT_SECONDARY, False, PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "Yuna (win/claude/opus/clawdbot)  |  Arin (wsl/claude/opus/clawdbot)  |  kingyy (win/vscode/opus/hum)  |  Sohee (win/antigravity/gemini-pro-high/clawdbot)",
             16, ACCENT_GREEN, False, PP_ALIGN.CENTER)


# ============================================================================
# Process entry by type
# ============================================================================

def process_entry(prs, date_str, author, entry):
    source_md = entry.get("source_md")
    t = entry["type"]

    # Special handling for table entries
    if source_md == "20260216_SAM3D_Pipeline.md":
        # SAM3D Pipeline deep dive with table and flowchart
        make_text_page(prs, date_str, author, "Analysis", ACCENT_PURPLE,
                       entry["title"], entry["summary"], entry.get("key_points", []))
        first_slide = prs.slides[len(prs.slides) - 1]
        make_pipeline_table_slide(prs, date_str, "SAM3D Pipeline: 6-Step Breakdown")
        make_pipeline_flowchart_slide(prs, date_str, "SAM3D Pipeline: Visual Flow")
    elif source_md == "20260224_SAM3D_TRELLIS_Architecture_Comparison.md":
        # Architecture comparison with table
        make_text_page(prs, date_str, author, "Analysis", ACCENT_PURPLE,
                       entry["title"], entry["summary"], entry.get("key_points", []))
        first_slide = prs.slides[len(prs.slides) - 1]
        make_architecture_table_slide(prs, date_str, "SAM3D vs TRELLIS v1 vs TRELLIS.2")
    elif source_md == "20260224_TRELLIS2_vs_SAM3D_Dining_Comparison.md":
        # TRELLIS2 vs SAM3D comparison with per-object GIFs
        make_text_page(prs, date_str, author, "Analysis", ACCENT_PURPLE,
                       entry["title"], entry["summary"], entry.get("key_points", []))
        first_slide = prs.slides[len(prs.slides) - 1]
        if entry.get("images"):
            make_large_image_pages(prs, date_str, entry["title"], entry["images"])
        if entry.get("rounds"):
            make_run_rounds(prs, date_str, entry["title"], entry["rounds"])
    elif t == "run":
        make_run_overview(prs, date_str, author, entry["title"], entry["summary"],
                          entry.get("input_img"), entry.get("output_img"))
        first_slide = prs.slides[len(prs.slides) - 1]
        if entry.get("rounds"):
            make_run_rounds(prs, date_str, entry["title"], entry["rounds"])
    elif t == "analysis":
        make_text_page(prs, date_str, author, "Analysis", ACCENT_PURPLE,
                       entry["title"], entry["summary"], entry.get("key_points", []))
        first_slide = prs.slides[len(prs.slides) - 1]
        if entry.get("images"):
            make_large_image_pages(prs, date_str, entry["title"], entry["images"])
    elif t == "summary":
        kp = entry.get("key_points", [])
        imgs = entry.get("images", [])
        if imgs:
            # Page 1: all text, Page 2+: images
            make_text_page(prs, date_str, author, "Summary", ACCENT_ORANGE,
                           entry["title"], entry["summary"], kp)
            first_slide = prs.slides[len(prs.slides) - 1]
            make_large_image_pages(prs, date_str, entry["title"], imgs)
        else:
            # No images: split key points across 2 pages
            mid = max(len(kp) // 2, 1)
            make_text_page(prs, date_str, author, "Summary", ACCENT_ORANGE,
                           entry["title"], entry["summary"], kp[:mid])
            first_slide = prs.slides[len(prs.slides) - 1]
            make_text_page(prs, date_str, author, "Summary (cont.)", ACCENT_ORANGE,
                           entry["title"], "", kp[mid:] if len(kp) > mid else ["(continued)"])
    else:
        first_slide = None

    if source_md and first_slide:
        url = GITHUB_DOCS_BASE + source_md
        add_hyperlink(first_slide, Inches(8.5), Inches(7.15), Inches(4.5), Inches(0.3),
                      "View on GitHub \u2192", url, 10, ACCENT_BLUE)


# ============================================================================
# DATA — organized by date (newest first), each date has entries
# ============================================================================

DATES = [
    # -------------------------------------------------------------------------
    # 2026-02-24
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-24",
        "author": "kingyy (gcp-L4/vscode/opus/hum)",
        "entries": [
            {
                "type": "analysis",
                "title": "TRELLIS2 vs SAM3D — Dining Scene Comparison",
                "source_md": "20260224_TRELLIS2_vs_SAM3D_Dining_Comparison.md",
                "summary": "Head-to-head comparison: SAM3D 2.3x better IoU (0.53 vs 0.23), 2.1x faster, "
                           "28x smaller GLBs. Gap driven by missing SS pose prior in TRELLIS2.",
                "key_points": [
                    "SAM3D avg IoU 0.53 (best: neck_pillow 0.90) vs TRELLIS2 0.23 (best: pillow_and_blanket 0.58)",
                    "SAM3D total 568s (9.5 min) vs TRELLIS2 1186s (19.8 min) — 2.1x faster",
                    "SAM3D GLBs 1.2-1.7 MB (80K verts) vs T2 18-81 MB (0.5-2.4M verts); more polys != better alignment",
                    "Only 1 comparable object improved with T2 (sofa_cover +0.15); newspaper regressed most (0.84 → 0.15)",
                    "Root cause: TRELLIS2 has no SS pose prior; Module 4b confirmed SS doesn't transfer cross-architecture",
                ],
                "images": [
                    ("T2 Masks", "test_results_images/trellis_comparison/t2_mask_grid.png"),
                    ("T2 vs Target", "test_results_images/trellis_comparison/t2_comparison_fixed.png"),
                ],
                "rounds": [
                    ("SAM3D: neck_pillow (IoU 0.90)", "../output/sam3d_dining_t1/rotation_gifs/neck_pillow_y_rotation.gif"),
                    ("SAM3D: newspaper (IoU 0.84)", "../output/sam3d_dining_t1/rotation_gifs/newspaper_y_rotation.gif"),
                    ("SAM3D: broken_tile (IoU 0.72)", "../output/sam3d_dining_t1/rotation_gifs/broken_tile_y_rotation.gif"),
                    ("SAM3D: placemat (IoU 0.62)", "../output/sam3d_dining_t1/rotation_gifs/placemat_y_rotation.gif"),
                    ("SAM3D: table_with_tablecloth (IoU 0.48)", "../output/sam3d_dining_t1/rotation_gifs/table_with_flower_tablecloth_y_rotation.gif"),
                    ("SAM3D: sofa (IoU 0.25)", "../output/sam3d_dining_t1/rotation_gifs/sofa_with_patterned_cover_y_rotation.gif"),
                    ("SAM3D: colander (IoU 0.24)", "../output/sam3d_dining_t1/rotation_gifs/metal_colander_y_rotation.gif"),
                    ("SAM3D: chair (IoU 0.23)", "../output/sam3d_dining_t1/rotation_gifs/wooden_chair_y_rotation.gif"),
                    ("T2: pillow_and_blanket (IoU 0.58)", "../output/sam3d_dining_t2/rotation_gifs/pillow_and_blanket_y_rotation.gif"),
                    ("T2: sofa_cover (IoU 0.40)", "../output/sam3d_dining_t2/rotation_gifs/sofa_cover_y_rotation.gif"),
                    ("T2: chair (IoU 0.24)", "../output/sam3d_dining_t2/rotation_gifs/chair_y_rotation.gif"),
                    ("T2: tablecloth (IoU 0.20)", "../output/sam3d_dining_t2/rotation_gifs/tablecloth_y_rotation.gif"),
                    ("T2: newspaper (IoU 0.15)", "../output/sam3d_dining_t2/rotation_gifs/newspaper_y_rotation.gif"),
                    ("T2: chair_cover (IoU 0.10)", "../output/sam3d_dining_t2/rotation_gifs/chair_cover_y_rotation.gif"),
                    ("T2: plant (IoU 0.09)", "../output/sam3d_dining_t2/rotation_gifs/plant_y_rotation.gif"),
                    ("T2: pot_and_trivet (IoU 0.08)", "../output/sam3d_dining_t2/rotation_gifs/pot_and_trivet_y_rotation.gif"),
                ],
            },
            {
                "type": "analysis",
                "title": "SAM3D / TRELLIS Architecture Comparison — Code Lineage & Module 4b",
                "source_md": "20260224_SAM3D_TRELLIS_Architecture_Comparison.md",
                "summary": "SAM3D backbone derived from TRELLIS v1 (MIT license) with Meta's MoT pose heads. "
                           "TRELLIS.2 uses O-Voxel octree (4B params). Module 4b: SS pose doesn't transfer.",
                "key_points": [
                    "SAM3D = TRELLIS v1 backbone + Meta additions (MoT attention, pointmap conditioning, pose decoder)",
                    "Code lineage: identical class names (SparseStructureFlowModel, SLatFlowModel, SLatGaussianDecoder), "
                    "identical imports, identical __init__ signatures — strong evidence of TRELLIS v1 fork",
                    "SAM3D ships BOTH original (T1-identical) and MoT-upgraded files (mot_sparse_structure_flow.py)",
                    "TRELLIS v1: SLAT dense 64^3, Structured Latent, ~500M params, MIT license (CVPR'25 Spotlight)",
                    "TRELLIS.2: O-Voxel octree up to 1536^3, 3-stage DiT, ~4B params, flexible resolution",
                    "SAM3D: 7 models ~3-4GB VRAM; TRELLIS.2: ~9-10GB VRAM (4B params)",
                    "Key difference: SAM3D's SS model predicts object pose (S,R,T) end-to-end — TRELLIS.2 has no equivalent",
                    "Module 4b (negative result): SS model's pose predictions trained with T1 decoder output space, "
                    "don't transfer to T2 mesh geometry — TRELLIS2+SS IoU = 0.157, same as baseline 0.157",
                    "Implication: TRELLIS.2 needs its own pose estimation module or external pose pipeline to match SAM3D alignment",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-23
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-23",
        "author": "kingyy (gcp-L4/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Modular Pipeline v4 — Gemini VLM + SS Pose Fix",
                "source_md": "20260223_Modular_Pipeline_v4_Gemini_All_Masks.md",
                "summary": "Switched VLM to Gemini 2.5 Flash, kept all 10 masks. SS pose fix (commit e93d61a) "
                           "improved avg IoU +0.055 (0.396 → 0.451). 9/10 objects in final scene.",
                "input_img": "test_results_images/modular_dining_v4/recognize_grid.png",
                "output_img": "test_results_images/modular_dining_v4_ss/side_by_side.png",
                "rounds": [],
                "key_points": [
                    "VLM: GPT-4o → Gemini 2.5 Flash (free tier, 5 req/min with backoff)",
                    "All 10 masks kept (no background rejection); 9/10 objects in final scene",
                    "SS pose fix: registration passes checkpoint_path → ICP gets good starting pose",
                    "Biggest improvements: placemat +0.18, plant +0.12, neck_pillow +0.10, armchair +0.13",
                    "IoU filter: exclude objects < 0.15 from render to prevent visual corruption",
                ],
                "images": [
                    ("Scene Render", "test_results_images/modular_dining_v4_ss/scene_render.png"),
                    ("Flat Render", "test_results_images/modular_dining_v4_ss/flat_scene_render.png"),
                    ("Projection Overlay", "test_results_images/modular_dining_v4_ss/projection_overlay.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "Modular Pipeline v3 — Large Mask Filter Test",
                "source_md": "20260223_Modular_Pipeline_v3_Large_Masks.md",
                "summary": "Raised max_area_ratio 50% → 95% to capture large masks. 2 new large masks detected "
                           "but GPT-4o correctly rejected both as background. No net change in scene.",
                "key_points": [
                    "Segment filter: max_area_ratio 0.50 → 0.95; 10 masks (vs 9 in v2)",
                    "Both large masks (33.3%, 17.0%) classified as background by GPT-4o",
                    "Same 8 foreground objects reconstructed; IoU variance from non-deterministic VLM naming",
                    "Finding: large mask filter has no practical impact; VLM correctly identifies background",
                ],
                "images": [
                    ("v3 Masks", "test_results_images/modular_dining_v3/scene_render.png"),
                    ("v3 Side-by-Side", "test_results_images/modular_dining_v3/flat_side_by_side.png"),
                ],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-22
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-22",
        "author": "kingyy (gcp-L4/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "SAM3D TRELLIS1 Dining — Full Pipeline Results",
                "source_md": "20260222_TRELLIS1_Dining_Results.md",
                "summary": "Complete SAM3D pipeline: 8 objects, avg IoU 0.53, 9.5 min total. "
                           "Batch TRELLIS1 with model caching. Best: neck_pillow 0.90.",
                "input_img": "test_results_images/modular_dining_v2/target.jpg",
                "output_img": "test_results_images/modular_dining_v2/registration_v2_side_by_side.png",
                "rounds": [],
                "key_points": [
                    "8 objects: neck_pillow (0.90), newspaper (0.84), broken_tile (0.72), placemat (0.62)",
                    "table_with_tablecloth (0.48), sofa (0.25), colander (0.24), chair (0.23)",
                    "SAM 52.7s + TRELLIS1 batch 515.2s = 567.9s (9.5 min) total",
                    "Rotation GIFs: Y-axis turntable + X-axis tumble per object via Blender Cycles 512x512",
                ],
                "images": [
                    ("Scene Comparison", "test_results_images/modular_dining_v2/registration_v2_projection_overlay.png"),
                ],
            },
            {
                "type": "run",
                "title": "Modular SAM3D Pipeline — 5-Module Design",
                "source_md": "20260222_Modular_Pipeline_Dining_Results.md",
                "summary": "Decomposed monolithic pipeline into 5 modules (segment→recognize→monodepth→reconstruct→register). "
                           "v2 fix: SS pose from checkpoint → avg IoU 0.414 → 0.543.",
                "input_img": "test_results_images/modular_dining_v2/segment_all_masks_grid.png",
                "output_img": "test_results_images/modular_dining_v2/registration_v2_scene_render.png",
                "rounds": [],
                "key_points": [
                    "5 modules: segment (SAM) → recognize (VLM) → monodepth (MoGe) → reconstruct (TRELLIS) → register (ICP+grad)",
                    "JSON manifests between modules for independence and replay",
                    "v2 critical fix: loads SS pose from checkpoint NPZ, applies square padding + isotropic intrinsics",
                    "Timing: 30s SAM + 45s VLM + 10s MoGe + 490s TRELLIS + 85s registration = ~660s (11 min)",
                ],
                "images": [
                    ("v2 Flat Render", "test_results_images/modular_dining_v2/registration_v2_flat_scene_render.png"),
                ],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-21
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-21",
        "author": "kingyy (gcp-L4/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "SAM3D Batch Worker — 2x Speedup via Model Caching",
                "source_md": "20260221_SAM3D_Batch_Worker_Results.md",
                "summary": "Batch worker loads TRELLIS once, processes all objects sequentially. "
                           "8 objects in 491s (8.2 min) vs baseline 1002s (16.7 min). Flash-attn only 2.4% faster.",
                "input_img": "test_results_images/modular_dining_v2/target.jpg",
                "output_img": "test_results_images/modular_dining_v2/registration_v2_scene_render.png",
                "rounds": [],
                "key_points": [
                    "Batch mode: load TRELLIS once, process N objects — eliminates ~27s/object model reload",
                    "491s (8.2 min) for 8 objects vs 1002s (16.7 min) baseline — 2x speedup",
                    "Flash-attention only 2.4% faster (attention is ~24% of pipeline; decoder dominates)",
                    "Zero OOM failures in batch mode vs 1 failure in baseline",
                    "GPU reload fix: after post-opt offloads to CPU, reload_pipeline_to_gpu() (2.9s each)",
                ],
                "images": [],
            },
            {
                "type": "analysis",
                "title": "TRELLIS2 vs TRELLIS1 — Initial Comparison",
                "source_md": "20260221_TRELLIS2_vs_TRELLIS1_Comparison.md",
                "summary": "TRELLIS2 produces visually superior PBR meshes but 2.8x slower. "
                           "IoU 0.22 (T2) vs 0.51 (T1) due to high-poly density mismatch in ICP.",
                "key_points": [
                    "TRELLIS2: 500K-2.4M vertices with PBR textures; TRELLIS1: 3K-13K vertices",
                    "T2 2.8x slower: 1409s vs 491s (model load 169s, separate subprocess required)",
                    "IoU 0.22 (T2) vs 0.51 (T1) — mesh density mismatch in ICP point sampling",
                    "Both use identical layout_post_optimization; difference is in mesh quality/density",
                    "Potential fix: decimate T2 meshes before ICP alignment",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-19
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-19",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "SAM3D Dining — ICP Pose Refinement + Overlay",
                "source_md": "20260219_SAM3D_Dining_ICP_Overlay.md",
                "summary": "ICP alignment on all 9 dining GLBs with v9 mask growth. "
                           "4/9 objects accepted ICP improvement. Scene overlay rendered.",
                "input_img": "test_results_images/sam3d_dining_icp/scene_photo_comparison.png",
                "output_img": "test_results_images/sam3d_dining_icp/scene_2d_comparison.png",
                "rounds": [],
                "key_points": [
                    "4/9 objects accepted ICP: chair_cushion, placemat, sofa (+0.3%), travel_pillow",
                    "Well-aligned objects (newspaper 1.5%, placemat 0.5%) already optimal — no ICP improvement",
                    "round_table/wooden_chair: high depth error from geometry (draping cloth, sparse frame), not pose",
                    "Open3D point-to-point ICP (max_correspondence=0.3m) with depth-scale correction",
                ],
                "images": [
                    ("Scene Depth", "test_results_images/sam3d_dining_icp/scene_depth_after.png"),
                ],
            },
            {
                "type": "run",
                "title": "wooden_chair — ICP Pipeline Evolution (4 Runs)",
                "source_md": "20260219_wooden_chair_ICP_pipeline_comparison.md",
                "summary": "Four sequential runs: --scene-image no effect, coarse ICP (0.1m) achieves 3x improvement "
                           "(0.09 → 0.26 IoU), Adam rejection bug fixed to preserve ICP gains.",
                "input_img": "test_results_images/20260219_wooden_chair_icp/r2_glb_render.png",
                "output_img": "test_results_images/20260219_wooden_chair_icp/r4_icp_preserve_proj.png",
                "rounds": [],
                "key_points": [
                    "R1 baseline (0.17), R2 +scene-image (0.17) — same pose, MoGe not the limiting factor",
                    "R3 two-pass ICP (0.1m→0.05m): coarse pass finds correspondences, IoU 0.27",
                    "R4: preserve ICP on Adam fail → IoU 0.26, chair at correct location",
                    "Ceiling ~0.27: TRELLIS initial pose tilted ~45°, Adam can't escape local minimum",
                ],
                "images": [
                    ("3-Way Comparison", "test_results_images/20260219_wooden_chair_icp/comparison_3way.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "NaN Root Cause + Pipeline Step 5 Corrections",
                "source_md": "20260219_NaN_Root_Cause_Experiment.md",
                "summary": "Controlled experiment confirms scene-image fix working (sparse NaN vs all NaN). "
                           "Pipeline corrections: mask growth before ICP, Adam always runs, depth filter clarified.",
                "key_points": [
                    "Per-object MoGe: all-NaN depth (4.7% mask coverage); scene MoGe: sparse-NaN (valid interior)",
                    "grow_mask_v9 stalls with NaN: depth comparisons return False, blocking growth",
                    "Correction: Stage 5a uses grown mask, not raw SAM mask",
                    "Correction: Adam always runs regardless of ICP outcome",
                    "wooden_chair: 4,685 SLAT voxels (sparse) vs 21,608 for keyboard — intrinsically hard",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-18
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-18",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "analysis",
                "title": "Mask Growth Algorithm Exploration (6 Variants)",
                "source_md": "20260218_SAM3D_Dining_v9_RayCast_MaskGrowth.md",
                "summary": "Systematic comparison of 6 mask growth algorithms on dining scene (v5-v9): "
                           "normal-consistency, plane-distance, RANSAC, depth gates, ray-cast. "
                           "Ray-cast (v9) best for convex objects, plane-distance best overall.",
                "key_points": [
                    "v5 Normal-Consistency: adaptive 15-60° threshold; flat objects conservative, curved objects hit cap",
                    "v6 Plane-Distance (EDT): local reference fixes curved surfaces (chair_cushion +6→+6,436 px); best overall",
                    "v7 RANSAC: 8-sector plane fit; over-conservative on flat surfaces (5mm floor collapses)",
                    "v8 RANSAC v2: relaxed 3cm floor doubles pixels but still 50% below plane-distance",
                    "v9 Ray-Cast: 8-direction first-hit depth; round_table +44K, sofa +12.7K — best for convex/filled objects",
                    "Key insight: hull pixels surrounded by mask pixels in all directions → rays capture full local geometry",
                ],
                "images": [
                    ("v5 Normal (round_table)", "test_results_images/sam3d_dining_v5/round_table_with_tablecloth_mask_growth.png"),
                    ("v6 Plane-Dist (chair_cushion)", "test_results_images/sam3d_dining_plane_dist/chair_cushion_mask_growth.png"),
                    ("v9 Ray-Cast (round_table)", "test_results_images/sam3d_dining_v9/round_table_with_tablecloth_mask_growth.png"),
                    ("v9 Ray-Cast (chair_legs)", "test_results_images/sam3d_dining_v9/chair_legs_mask_growth.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "Convex Hull Results + Normal-Consistency on GreenTea",
                "source_md": "20260218_SAM3D_ConvexHull_Results.md",
                "summary": "Scene-image MoGe is a VIGA modification (not original SAM3D). "
                           "Normal-consistency replaces Sobel for mask growth; handles shadows better, limited on curves.",
                "key_points": [
                    "Scene-image MoGe is our modification: original SAM3D runs MoGe on per-object masked RGBA",
                    "Original design fails when objects <10% of image (masked image mostly black → MoGe NaN)",
                    "Normal method correctly rejects shadow-contaminated regions (green_tea_bottle)",
                    "Curved surface limitation: global mean reference breaks for cylinders (normal angle 0-90°)",
                    "Per-pixel local reference approach suggested as future improvement",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-17
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-17",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Convex Hull + Scene-Image MoGe — GreenTea Results",
                "source_md": "20260217_SAM3D_ConvexHull_GreenTea_Results.md",
                "summary": "Scene-image MoGe fix recovered 4/5 objects from all-NaN failure. "
                           "IoU 0.45-0.95 (ito_en_bottle 0.95, envelope 0.86, headphones 0.83). "
                           "Green_tea_bottle degenerated due to shadow contamination.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/greentea/target.png",
                "rounds": [],
                "key_points": [
                    "Scene-image MoGe bypasses NaN failures in per-object masked images (<30% visible pixels)",
                    "4/5 objects: ito_en_bottle (0.95), envelope (0.86), headphones (0.83), keyboard (0.65)",
                    "green_tea_bottle degenerated to flat disk — shadow contamination in SAM mask",
                    "Convex hull mask growth stops at depth edges, gives TRELLIS more silhouette context",
                    "~70-80 min end-to-end for 5 objects including Blender GIF rendering",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-16
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-16",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "SAM3D Dining v2/v3 — Post-Transform-Fix Full Run",
                "source_md": "20260216_SAM3D_Dining_v2_Results.md",
                "summary": "First 9-object dining run after Transform3d fix. Scale correction + Sobel mask growth. "
                           "chair_legs depth error 55% → 5.3%. Added scene-level 2D comparison viz.",
                "input_img": "test_results_images/sam3d_dining_v2/full_scene_comparison.png",
                "output_img": "test_results_images/sam3d_dining_v3/scene_2d_comparison.png",
                "rounds": [],
                "key_points": [
                    "3 transform fixes applied: translation row, pre-transform sign, dead post-transforms removed",
                    "Depth alignment improved: chair_legs 55% → 5.3%, sofa 51% → 7.8%",
                    "Best: placemat (0.6% error), newspaper (1.6%), chair_cushion (2.7%)",
                    "v3: added scene-level 2D comparison — SAM masks (left) vs GLB projections (right)",
                ],
                "images": [
                    ("v2 Scene Compare", "test_results_images/sam3d_dining_v2/full_scene_comparison.png"),
                    ("v3 2D Compare", "test_results_images/sam3d_dining_v3/scene_2d_comparison.png"),
                    ("Depth Dashboard", "test_results_images/sam3d_dining_v2/depth_diagnostic_dashboard.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "Depth Alignment Analysis",
                "source_md": "20260216_Depth_Alignment_Analysis.md",
                "summary": "Diagnostic of depth alignment: layout optimizer prioritizes 2D silhouette, not depth.",
                "key_points": [
                    "Test: per-vertex depth well-aligned for flat objects (1.4-5.3%) but severely misaligned for complex shapes (25-55%)",
                    "Root cause: layout post-optimization (Stage 5c) optimizes only 2D silhouette IoU, no depth loss",
                    "Recommended fix: add depth loss term to differentiable rendering in Stage 5c",
                ],
                "images": [],
            },
            {
                "type": "analysis",
                "title": "SAM3D Pipeline Deep Dive — 6-Step Architecture",
                "source_md": "20260216_SAM3D_Pipeline.md",
                "summary": "Detailed decomposition of the full SAM3D pipeline: SAM segmentation → MoGe depth "
                           "→ TRELLIS 3D reconstruction → pose decoding → layout optimization → GLB export. "
                           "Key insight: TRELLIS sees only RGBA image, MoGe pointmap used only for pose.",
                "key_points": [
                    "Step 1 (SAM ViT-H): Full scene → N binary masks; conda env 'sam' (Py3.10)",
                    "Step 2 (MoGe): Scene image → (H,W,3) pointmap + intrinsics; NOT a depth map — full 3D camera-space coords per pixel",
                    "Step 3 (TRELLIS): Per-object masked RGBA → 3D mesh; SS (2 steps) + SLAT (12 steps) + dual decoder (32 Gaussians/voxel)",
                    "Critical: TRELLIS receives ONLY RGBA image, NOT MoGe depth — shape comes purely from image appearance",
                    "Step 4 (Pose Decode): SS model predicts initial S, R, T from sparse structure + pointmap conditioning",
                    "Step 5 (Layout Optimization): 3 sub-stages — (5a) pointmap-based coarse, (5b) ICP refinement, (5c) differentiable silhouette rendering",
                    "MoGe intrinsics are single source of camera calibration — wrong intrinsics cascade to all downstream alignment",
                    "scene-image fix: original per-object SSI normalization fails at <5% pixel coverage → bypasses SSI with raw MoGe pointmap",
                    "Step 6 (Export): Bake S, R, T into vertex positions → GLB with vertex colors in PyTorch3D camera space",
                    "Coordinate systems: MoGe camera (X-right, Y-down, Z-forward) → PyTorch3D (X-left, Y-up, Z-forward) → Blender GLTF (Y-up → Z-up)",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-15
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-15",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "SAM3D Dining Scene — First Full Pipeline Run",
                "source_md": "20260215_SAM3D_Dining_Scene_Results.md",
                "summary": "End-to-end pipeline on dining scene: 9 objects segmented, reconstructed, and rendered. "
                           "103 min total. 8/9 batch + 1 OOM rerun. ~89 MB GLBs.",
                "input_img": "test_results_images/sam3d_dining_v2/full_scene_comparison.png",
                "output_img": "test_results_images/sam3d_dining_v2/full_scene_comparison.png",
                "rounds": [],
                "key_points": [
                    "Input: 771x1024 dining scene (resized from 3072x4080 to avoid SAM CUDA OOM)",
                    "9 objects: chair, cushion, legs, table, sofa, pillow, newspaper, strainer, placemat",
                    "8/9 batch completed; chair_legs OOM'd and required solo rerun (15 min, 12.5K sparse coords)",
                    "MoGe depth 0.70-3.01m, camera intrinsics fx/fy=701.1px, cx=385.5, cy=512.0",
                ],
                "images": [],
            },
            {
                "type": "analysis",
                "title": "SAM3D Transform3d Alignment Fix",
                "source_md": "20260215_SAM3D_Alignment_Fix.md",
                "summary": "Fixed three critical bugs in Transform3d wrapper: translation in wrong row, "
                           "negated pre-transform X-axis, dead post-transforms. 6/6 greentea objects now aligned.",
                "key_points": [
                    "Bug 1: Translation in column 3 (wrong) vs row 3 (correct) for PyTorch3D row-vector convention",
                    "Bug 2: Pre-transform had negated X-axis, mirroring all meshes horizontally",
                    "Bug 3: Three post-transforms canceled to identity — removed as dead code",
                    "Per-object pixel accuracy 0.4-2.5% except large flat surfaces (desk 10.4%)",
                    "Render axis fixes: vertical flip (Y-down→Y-up), horizontal flip (X-left vs X-right)",
                ],
                "images": [],
            },
            {
                "type": "analysis",
                "title": "SAM3D / TRELLIS Deep Dive Study",
                "source_md": "20260215_SAM3D_TRELLIS_Study.md",
                "summary": "Architecture study: SAM ViT-H → TRELLIS reconstruction. Bottleneck is dual decoder (50%). "
                           "TRELLIS2 (4B params) promises ~3s/object on H100.",
                "key_points": [
                    "Pipeline: SAM segmentation → TRELLIS 3D (SS + SLAT + decoder)",
                    "Bottleneck: dual decoder 50%, SLAT sampling 30%, model loading 5%",
                    "TRELLIS2 (Dec 2025): 4B params, ~3s/object on H100 but consumer GPU VRAM unknown",
                    "Optimization roadmap: reduce SLAT steps, cache model, parallelize objects",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-12
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-12",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "analysis",
                "title": "VIGA Pipeline Revisited",
                "source_md": "20260212_VIGA_PIPELINE_REVISITED.md",
                "summary": "Detailed pipeline architecture: SAM3D reconstruction, Meshy text-to-3D, "
                           "Generator/Verifier agents, Blender executor. Per-module inputs/outputs, "
                           "pipeline modes, environment setup, and agent tools reference.",
                "key_points": [
                    "Pipeline: Target Image -> SAM3D (image->3D) + Meshy (text->3D) -> Generator (Blender Python) -> Executor (EEVEE render) -> Verifier (feedback) -> loop",
                    "SAM3D: SAM segmentation (ViT-H) -> per-object binary masks -> SAM3D/TRELLIS reconstruction -> GLB with vertex colors + transform JSON",
                    "Meshy: local cache fuzzy match (instant) -> API fallback (text-to-3D preview -> refine -> download GLB, ~7 min/object)",
                    "Generator: VLM (GPT-5) writes complete Blender Python scripts per round, calls tools (execute_and_evaluate, get_better_object, initialize_plan)",
                    "Verifier: VLM compares render to target, provides structured text feedback, can inspect scene from multiple angles",
                    "3 pipeline modes: Meshy-only (get_asset), SAM3D+Meshy (get_asset_sam3d), Pre-computed SAM3D (--sam3d-results)",
                    "SAM3D vs Meshy: SAM3D captures scene layout + good for simple shapes; Meshy gives UV-textured meshes for complex objects",
                    "Environment: agent env (Python 3.10, openai, mcp) + sam3d_viga env (Python 3.11, torch 2.5.1+cu121, kaolin 0.17.0)",
                    "Hardware: RTX 5080 16GB, Blender 4.5.5 LTS, Windows 11",
                    "Output per round: scripts/{N}.py + renders/{N}/Camera.png + renders/{N}/state.blend",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-11
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-11",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Dynamic Scene: Artist Run 2 (SAM3D Fixed)",
                "source_md": "20260211_DynamicScene_Artist_Run2_Results.md",
                "summary": "SAM3D injection bug fixed. Pre-computed GLBs loaded via --sam3d-results. "
                           "Generator kept 5/8 SAM3D objects (fruits), replaced jug/plate/pear with Meshy cache. "
                           "50% faster than Run 1 (99 min vs 203 min). 21 rounds.",
                "input_img": "test_results_images/dynamic_artist_run2/target.png",
                "output_img": "test_results_images/dynamic_artist_run2/gifs/round_19_rotation.gif",
                "rounds": [
                    ("R2 First SAM3D Import", "test_results_images/dynamic_artist_run2/gifs/round_2.gif"),
                    ("R4 Objects on Table", "test_results_images/dynamic_artist_run2/gifs/round_4.gif"),
                    ("R5 SAM3D+Meshy Mix", "test_results_images/dynamic_artist_run2/gifs/round_5.gif"),
                    ("R15 Best Static", "test_results_images/dynamic_artist_run2/gifs/round_15.gif"),
                    ("R16 Ball Added", "test_results_images/dynamic_artist_run2/gifs/round_16.gif"),
                    ("R17 Ball Impact (360)", "test_results_images/dynamic_artist_run2/gifs/round_17.gif"),
                    ("R19 Rotation (360)", "test_results_images/dynamic_artist_run2/gifs/round_19_rotation.gif"),
                    ("R19 Animation (180f)", "test_results_images/dynamic_artist_run2/gifs/round_19_animation.gif"),
                ],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-10
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-10",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Dynamic Scene: Artist Run 1 (SAM3D + Meshy)",
                "source_md": "20260210_DynamicScene_Artist_Run1_Results.md",
                "summary": "First SAM3D-initialized dynamic scene. 8 objects reconstructed but injection bug "
                           "(generator.py:81) meant Generator never received paths — replaced ALL with Meshy. "
                           "Ball physics too subtle. 18 rounds, 203 min.",
                "input_img": "test_results_images/dynamic_artist_run1/target.png",
                "output_img": "test_results_images/dynamic_artist_run1/gifs/round_18.gif",
                "rounds": [
                    ("SAM3D: Ceramic Jug", "test_results_images/dynamic_artist_run1/sam3d_renders/ceramic_jug.gif"),
                    ("SAM3D: Green Pears", "test_results_images/dynamic_artist_run1/sam3d_renders/green_pears.gif"),
                    ("SAM3D: Orange Pears", "test_results_images/dynamic_artist_run1/sam3d_renders/orange_pears.gif"),
                    ("SAM3D: Plate+Fruits", "test_results_images/dynamic_artist_run1/sam3d_renders/plate_with_fruits.gif"),
                    ("SAM3D: Orange Pear", "test_results_images/dynamic_artist_run1/sam3d_renders/orange_pear.gif"),
                    ("SAM3D: Green Apple", "test_results_images/dynamic_artist_run1/sam3d_renders/green_apple.gif"),
                    ("SAM3D: Green Pear", "test_results_images/dynamic_artist_run1/sam3d_renders/green_pear.gif"),
                    ("SAM3D: Green Apple 2", "test_results_images/dynamic_artist_run1/sam3d_renders/green_apple_1.gif"),
                    ("R2 First Render", "test_results_images/dynamic_artist_run1/gifs/round_2.gif"),
                    ("R9 Scene Refinement", "test_results_images/dynamic_artist_run1/gifs/round_9.gif"),
                    ("R14 Best Static", "test_results_images/dynamic_artist_run1/gifs/round_14.gif"),
                    ("R16 Ball Added (360)", "test_results_images/dynamic_artist_run1/gifs/round_16.gif"),
                    ("R16 Frame 1", "test_results_images/dynamic_artist_run1/round_16/Camera_f0001.png"),
                    ("R16 Frame 320", "test_results_images/dynamic_artist_run1/round_16/Camera_f0320.png"),
                    ("R18 Final (360)", "test_results_images/dynamic_artist_run1/gifs/round_18.gif"),
                    ("R18 Frame 1", "test_results_images/dynamic_artist_run1/round_18/Camera_f0001.png"),
                    ("R18 Frame 260", "test_results_images/dynamic_artist_run1/round_18/Camera_f0260.png"),
                    ("R18 Frame 520", "test_results_images/dynamic_artist_run1/round_18/Camera_f0520.png"),
                ],
            },
            {
                "type": "run",
                "title": "SAM3D+Meshy Static Scene Run 1",
                "source_md": "20260210_SAM3D_Meshy_Run1_Results.md",
                "summary": "First combined SAM3D+Meshy static scene. SAM3D reconstructed 6 objects but Generator "
                           "replaced all with Meshy (4 cached, 1 new API). 25 rounds, 2h14m.",
                "input_img": "test_results_images/sam3d_meshy_run1/target.png",
                "output_img": "test_results_images/sam3d_meshy_run1/round_17_Cam.png",
                "rounds": [
                    ("R2", "test_results_images/sam3d_meshy_run1/round_02_Camera.png"),
                    ("R5", "test_results_images/sam3d_meshy_run1/round_05_Camera.png"),
                    ("R6", "test_results_images/sam3d_meshy_run1/round_06_Camera.png"),
                    ("R8", "test_results_images/sam3d_meshy_run1/round_08_Cam.png"),
                    ("R10", "test_results_images/sam3d_meshy_run1/round_10_Cam.png"),
                    ("R11", "test_results_images/sam3d_meshy_run1/round_11_Cam.png"),
                    ("R12", "test_results_images/sam3d_meshy_run1/round_12_Cam.png"),
                    ("R13", "test_results_images/sam3d_meshy_run1/round_13_Cam.png"),
                    ("R14", "test_results_images/sam3d_meshy_run1/round_14_Cam.png"),
                    ("R15", "test_results_images/sam3d_meshy_run1/round_15_Cam.png"),
                    ("R16", "test_results_images/sam3d_meshy_run1/round_16_Cam.png"),
                    ("R17 Final", "test_results_images/sam3d_meshy_run1/round_17_Cam.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "SAM3D+Meshy Combined Pipeline Design",
                "source_md": "20260210_SAM3D_Meshy_Combined_Pipeline.md",
                "summary": "Technical design for combining SAM3D reconstruction with Meshy text-to-3D. "
                           "SAM3D initializes scene from target image, Meshy replaces poor-quality objects.",
                "key_points": [
                    "Combined flow: SAM3D auto-init -> segment target -> reconstruct GLBs",
                    "Generator receives SAM3D object paths in memory context",
                    "Phase 1: Render SAM3D scene, evaluate each object quality",
                    "Phase 1: Call get_better_object (Meshy) for poor objects only",
                    "Phase 2: Compose scene with best GLBs (SAM3D or Meshy mix)",
                    "",
                    "Code changes to generator.py: capture SAM3D results, inject into memory",
                    "New prompt: get_asset_sam3d (three-phase workflow)",
                    "Bug fixes: kaolin DLL error, pipeline.yaml path resolution",
                    "Conda env mapping: sam env for SAM3D, agent env for Generator",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-09
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-09",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Static Scene: Green Tea Run 11",
                "source_md": "20260209_StaticScene_GreenTea_Run11_Results.md",
                "summary": "19 rounds with AABB collision avoidance and non-uniform table scaling. "
                           "Bottle Japanese label text visible by round 11. Cached Meshy assets.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/static_scene/20260209_143141/rotation_gif/round_19.gif",
                "rounds": [
                    ("R1", "test_results_images/static_scene/20260209_143141/rotation_gif/round_1.gif"),
                    ("R3", "test_results_images/static_scene/20260209_143141/rotation_gif/round_3.gif"),
                    ("R13", "test_results_images/static_scene/20260209_143141/rotation_gif/round_13.gif"),
                    ("R19 Final", "test_results_images/static_scene/20260209_143141/rotation_gif/round_19.gif"),
                ],
            },
            {
                "type": "run",
                "title": "Static Scene: Green Tea Run 10",
                "source_md": "20260209_StaticScene_GreenTea_Run10_Results.md",
                "summary": "19 rounds. Night-desk aesthetic with blue monitor glow. Introduced orient_group_min_z() "
                           "rotation helper for finding flattest object orientation. Cached Meshy assets.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": None,
                "rounds": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-08
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-08",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Static Scene: Green Tea Run 9 (Meshy API)",
                "source_md": "20260208_StaticScene_GreenTea_Run9_Results.md",
                "summary": "First run with Meshy API text-to-3D for all 5 assets. Dramatically higher quality than "
                           "SAM3D: recognizable PET bottle, RGB keyboard, headphones. Cinematic hero shot.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/static_scene/20260208_224146/rotation_gif/round_8.gif",
                "rounds": [
                    ("R3", "test_results_images/static_scene/20260208_224146/renders/round_3_Camera.png"),
                    ("R3 (360)", "test_results_images/static_scene/20260208_224146/rotation_gif/round_3.gif"),
                    ("R4", "test_results_images/static_scene/20260208_224146/renders/round_4_Camera.png"),
                    ("R5", "test_results_images/static_scene/20260208_224146/renders/round_5_Camera.png"),
                    ("R5 (360)", "test_results_images/static_scene/20260208_224146/rotation_gif/round_5.gif"),
                    ("R6", "test_results_images/static_scene/20260208_224146/renders/round_6_Camera.png"),
                    ("R7", "test_results_images/static_scene/20260208_224146/renders/round_7_Camera.png"),
                    ("R8 Final", "test_results_images/static_scene/20260208_224146/renders/round_8_Camera.png"),
                    ("R8 Final (360)", "test_results_images/static_scene/20260208_224146/rotation_gif/round_8.gif"),
                ],
            },
            {
                "type": "run",
                "title": "Static Scene: Green Tea Run 8",
                "source_md": "20260208_StaticScene_GreenTea_Run8_Results.md",
                "summary": "First run with Meshy API key configured (but never called — all 5 assets matched locally). "
                           "Slight regression in R6 (objects lost from camera view).",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/static_scene/20260208_204903/rotation_gif/round_9.gif",
                "rounds": [
                    ("R2", "test_results_images/static_scene/20260208_204903/renders/round_2_Camera.png"),
                    ("R2 (360)", "test_results_images/static_scene/20260208_204903/rotation_gif/round_2.gif"),
                    ("R4", "test_results_images/static_scene/20260208_204903/renders/round_4_Camera.png"),
                    ("R5", "test_results_images/static_scene/20260208_204903/renders/round_5_Camera.png"),
                    ("R5 (360)", "test_results_images/static_scene/20260208_204903/rotation_gif/round_5.gif"),
                    ("R6 Regression", "test_results_images/static_scene/20260208_204903/renders/round_6_Camera.png"),
                    ("R8", "test_results_images/static_scene/20260208_204903/renders/round_8_Camera.png"),
                    ("R9 Final", "test_results_images/static_scene/20260208_204903/renders/round_9_Camera.png"),
                    ("R9 Final (360)", "test_results_images/static_scene/20260208_204903/rotation_gif/round_9.gif"),
                ],
            },
            {
                "type": "run",
                "title": "Static Scene: Green Tea Run 7 (get_asset_simple)",
                "source_md": "20260208_StaticScene_GreenTea_Run7_Results.md",
                "summary": "New prompt setting forbids procedural geometry — only GLB imports allowed. "
                           "All 5 SAM3D assets loaded and arranged over 9 rounds. Zero procedural violations.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/static_scene/20260208_164519/rotation_gif/round_9.gif",
                "rounds": [
                    ("R1", "test_results_images/static_scene/20260208_164519/renders/round_1_Camera.png"),
                    ("R1 (360)", "test_results_images/static_scene/20260208_164519/rotation_gif/round_1.gif"),
                    ("R2", "test_results_images/static_scene/20260208_164519/renders/round_2_Camera.png"),
                    ("R3", "test_results_images/static_scene/20260208_164519/renders/round_3_Camera.png"),
                    ("R3 (360)", "test_results_images/static_scene/20260208_164519/rotation_gif/round_3.gif"),
                    ("R4", "test_results_images/static_scene/20260208_164519/renders/round_4_Camera.png"),
                    ("R5", "test_results_images/static_scene/20260208_164519/renders/round_5_Camera.png"),
                    ("R6", "test_results_images/static_scene/20260208_164519/renders/round_6_Camera.png"),
                    ("R6 (360)", "test_results_images/static_scene/20260208_164519/rotation_gif/round_6.gif"),
                    ("R8", "test_results_images/static_scene/20260208_164519/renders/round_8_Camera.png"),
                    ("R9 Final", "test_results_images/static_scene/20260208_164519/renders/round_9_Camera.png"),
                    ("R9 Final (360)", "test_results_images/static_scene/20260208_164519/rotation_gif/round_9.gif"),
                ],
            },
            {
                "type": "run",
                "title": "Static Scene: Green Tea Run 5 (First GLB Import)",
                "source_md": "20260208_StaticScene_GreenTea_Run5_Results.md",
                "summary": "First run where absolute path fix enabled successful GLB asset imports. "
                           "25 rounds, 14 renders. Progression from overexposed early renders to full scene.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/static_scene/20260208_050118/renders/round_20_Camera.png",
                "rounds": [
                    ("R3", "test_results_images/static_scene/20260208_050118/renders/round_3_Camera.png"),
                    ("R5", "test_results_images/static_scene/20260208_050118/renders/round_5_Camera.png"),
                    ("R7", "test_results_images/static_scene/20260208_050118/renders/round_7_Camera.png"),
                    ("R9", "test_results_images/static_scene/20260208_050118/renders/round_9_Camera.png"),
                    ("R11", "test_results_images/static_scene/20260208_050118/renders/round_11_Camera.png"),
                    ("R14", "test_results_images/static_scene/20260208_050118/renders/round_14_Camera.png"),
                    ("R16", "test_results_images/static_scene/20260208_050118/renders/round_16_Camera.png"),
                    ("R17", "test_results_images/static_scene/20260208_050118/renders/round_17_Camera.png"),
                    ("R19", "test_results_images/static_scene/20260208_050118/renders/round_19_Camera.png"),
                    ("R20 Final", "test_results_images/static_scene/20260208_050118/renders/round_20_Camera.png"),
                    ("CYCLES Final", "test_results_images/static_scene/20260208_050118/renders/Camera_cycles.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "Asset Pipeline Root Cause (Run 4)",
                "source_md": "20260208_StaticScene_GreenTea_AssetRun.md",
                "summary": "Root cause analysis: GLB imports failed silently due to relative path resolution in Blender. "
                           "Fix: meshy_api.py resolves to absolute paths with forward slashes.",
                "key_points": [
                    "Problem: Blender CWD != project root -> relative GLB paths fail silently",
                    "All 5 get_better_object calls returned valid paths but Blender couldn't find files",
                    "Generator fell back to procedural geometry for all objects",
                    "Fix in meshy_api.py: Path(previous_assets_dir).resolve() with forward slashes",
                    "Fix in prompts: instructions to use EXACT absolute file paths",
                    "Verified: Run 5 successfully imports all GLB assets after fix",
                ],
                "images": [
                    ("R4 (GLB import failed)", "test_results_images/static_scene/20260207_225603/renders/round_4_Camera.png"),
                    ("R8 (procedural fallback)", "test_results_images/static_scene/20260207_225603/renders/round_8_Camera.png"),
                    ("R13", "test_results_images/static_scene/20260207_225603/renders/round_13_Camera.png"),
                    ("R16", "test_results_images/static_scene/20260207_225603/renders/round_16_Camera.png"),
                    ("R18", "test_results_images/static_scene/20260207_225603/renders/round_18_Camera.png"),
                    ("R19 Final", "test_results_images/static_scene/20260207_225603/renders/round_19_Camera.png"),
                ],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-07
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-07",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Static Scene: Green Tea First Run (58 rounds)",
                "source_md": "20260207_StaticScene_GreenTea_Results.md",
                "summary": "First static scene run with --prompt-setting=none. 58 scripts, 36 renders, ~4 hours. "
                           "No GLB assets (MeshyAPI init failed). Severe scene drift after round 25.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/static_scene/20260207_030656/renders/round_58_Camera.png",
                "rounds": [
                    ("R7", "test_results_images/static_scene/20260207_030656/renders/round_7_Camera.png"),
                    ("R9", "test_results_images/static_scene/20260207_030656/renders/round_9_Camera.png"),
                    ("R13", "test_results_images/static_scene/20260207_030656/renders/round_13_Camera.png"),
                    ("R16", "test_results_images/static_scene/20260207_030656/renders/round_16_Camera.png"),
                    ("R19", "test_results_images/static_scene/20260207_030656/renders/round_19_Camera.png"),
                    ("R23", "test_results_images/static_scene/20260207_030656/renders/round_23_Camera.png"),
                    ("R25", "test_results_images/static_scene/20260207_030656/renders/round_25_Camera.png"),
                    ("R28", "test_results_images/static_scene/20260207_030656/renders/round_28_Camera.png"),
                    ("R30", "test_results_images/static_scene/20260207_030656/renders/round_30_Camera.png"),
                    ("R36", "test_results_images/static_scene/20260207_030656/renders/round_36_Camera.png"),
                    ("R42", "test_results_images/static_scene/20260207_030656/renders/round_42_Camera.png"),
                    ("R44", "test_results_images/static_scene/20260207_030656/renders/round_44_Camera.png"),
                    ("R47", "test_results_images/static_scene/20260207_030656/renders/round_47_Camera.png"),
                    ("R50", "test_results_images/static_scene/20260207_030656/renders/round_50_Camera.png"),
                    ("R52", "test_results_images/static_scene/20260207_030656/renders/round_52_Camera.png"),
                    ("R55", "test_results_images/static_scene/20260207_030656/renders/round_55_Camera.png"),
                    ("R57", "test_results_images/static_scene/20260207_030656/renders/round_57_Camera.png"),
                    ("R58 Final", "test_results_images/static_scene/20260207_030656/renders/round_58_Camera.png"),
                    ("CYCLES Final", "test_results_images/static_scene/20260207_030656/renders/CYCLES_final_Camera.png"),
                ],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-06
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-06",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Dynamic Scene: Artist Baseline",
                "source_md": "20260206_DynamicScene_Artist_Baseline_Results.md",
                "summary": "First dynamic scene on Windows with GPT-5. Cezanne still life. Procedural geometry only "
                           "(no SAM3D/Meshy). 12 rounds with physics-based ball-throwing animation.",
                "input_img": "test_results_images/dynamic_artist_baseline/target.png",
                "output_img": "test_results_images/dynamic_artist_baseline/gifs/round_12.gif",
                "rounds": [
                    ("R2", "test_results_images/dynamic_artist_baseline/gifs/round_2.gif"),
                    ("R5", "test_results_images/dynamic_artist_baseline/gifs/round_5.gif"),
                    ("R7", "test_results_images/dynamic_artist_baseline/gifs/round_7.gif"),
                    ("R9 (360)", "test_results_images/dynamic_artist_baseline/gifs/round_9.gif"),
                    ("R9 Frame 1", "test_results_images/dynamic_artist_baseline/keyframes/round_9/Camera_f0001.png"),
                    ("R9 Frame 90", "test_results_images/dynamic_artist_baseline/keyframes/round_9/Camera_f0090.png"),
                    ("R9 Frame 180", "test_results_images/dynamic_artist_baseline/keyframes/round_9/Camera_f0180.png"),
                    ("R11", "test_results_images/dynamic_artist_baseline/gifs/round_11.gif"),
                    ("R12 Final (360)", "test_results_images/dynamic_artist_baseline/gifs/round_12.gif"),
                    ("R12 Frame 1", "test_results_images/dynamic_artist_baseline/keyframes/round_12/Camera_f0001.png"),
                    ("R12 Frame 100", "test_results_images/dynamic_artist_baseline/keyframes/round_12/Camera_f0100.png"),
                    ("R12 Frame 200", "test_results_images/dynamic_artist_baseline/keyframes/round_12/Camera_f0200.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "Windows Compatibility Fixes",
                "source_md": "20260206_DynamicScene_Windows_Experience.md",
                "summary": "Four major Windows issues identified and fixed when running VIGA dynamic scene pipeline.",
                "key_points": [
                    "Issue 1: Blender path resolution",
                    "  Blender resolves relative paths to its own CWD, not project root",
                    "  Fix: .resolve() all Path objects before passing to Blender",
                    "",
                    "Issue 2: Windows path spaces",
                    "  Blender.exe path has spaces: 'C:\\Program Files\\Blender Foundation\\...'",
                    "  Fix: quote subprocess args containing spaces",
                    "",
                    "Issue 3: Subprocess pipe deadlock",
                    "  capture_output=True deadlocks on Windows with large outputs",
                    "  Fix: use temp files for stdout/stderr instead of pipes",
                    "",
                    "Issue 4: cp1252 encoding",
                    "  Windows defaults to cp1252, but scripts contain UTF-8 characters",
                    "  Fix: always specify encoding='utf-8' for file operations",
                    "",
                    "Files fixed: exec.py, investigator_core.py, generator.py, verifier.py, common.py",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-05
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-05",
        "author": "Yuna (win/claude/opus/clawdbot)",
        "entries": [
            {
                "type": "summary",
                "title": "SAM3D Pipeline Complete (5/5 Objects)",
                "source_md": "20260205_SAM3D_Pipeline_Complete.md",
                "summary": "End-to-end SAM + SAM3D pipeline completed. 100% success (5/5 objects) on green tea scene. "
                           "SAM segmentation -> mask extraction -> SAM3D reconstruction -> GLB export. ~2 hrs on RTX 5080.",
                "key_points": [
                    "Full pipeline: target.png -> SAM masks -> SAM3D 3D meshes -> GLB export -> Blender render -> GIF",
                    "5 objects: green tea bottle, Ito En bottle, keyboard, headphones, envelope",
                    "All objects successfully reconstructed with vertex colors",
                    "Both X-axis and Y-axis rotation GIFs generated per object",
                    "WebGL viewer integration for interactive 3D preview",
                    "Total processing: ~2 hours on NVIDIA RTX 5080 (16GB VRAM)",
                ],
                "images": [
                    ("Input Scene", "assets/01_greentea_input.jpg"),
                    ("Green Tea Bottle", "assets/green_tea_bottle_y_rotation.gif"),
                    ("Ito En Bottle", "assets/ito_en_bottle_y_rotation.gif"),
                    ("Alienware Keyboard", "assets/alienware_keyboard_y_rotation.gif"),
                    ("Headphones", "assets/headphones_y_rotation.gif"),
                    ("Envelope", "assets/envelope_y_rotation.gif"),
                ],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-04
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-04",
        "author": "Yuna (win/claude/opus/clawdbot)",
        "entries": [
            {
                "type": "summary",
                "title": "SAM Segmentation Results",
                "source_md": "20260204_SAM_SegmentationResults.md",
                "summary": "SAM segmented green tea desktop scene: 145 raw masks filtered to 8, identifying 5 objects.",
                "key_points": [
                    "Input: green tea desktop photograph",
                    "SAM produced 145 raw masks, filtered to 8 high-confidence regions",
                    "5 objects identified: green tea bottle, Ito En bottle, Alienware keyboard, headphones, envelope",
                    "Mask quality: clean edges, proper object boundaries",
                    "Ready for SAM3D batch 3D reconstruction",
                ],
                "images": [
                    ("Target Scene", "test_results_images/greentea/target.png"),
                    ("Green Tea Bottle Mask", "test_results_images/sam_init/green_tea_bottle.png"),
                    ("Ito En Bottle Mask", "test_results_images/sam_init/ito_en_bottle.png"),
                    ("Keyboard Mask", "test_results_images/sam_init/alienware_keyboard.png"),
                    ("Headphones Mask", "test_results_images/sam_init/headphones.png"),
                    ("Envelope Mask", "test_results_images/sam_init/envelope.png"),
                ],
            },
            {
                "type": "summary",
                "title": "SAM3D All-Masks Results",
                "source_md": "20260204_SAM3D_AllMasks.md",
                "summary": "SAM3D all-masks pipeline: 6 objects reconstructed with X and Y rotation GIFs.",
                "key_points": [
                    "6 objects: green tea bottle, Ito En bottle, bottle cap, label wrap, bottle neck, headphones",
                    "Both X-axis and Y-axis rotation GIFs generated for each object",
                    "Unexpected sub-object detection (bottle cap, label wrap, neck)",
                    "GLB generation statistics and mask coverage analysis",
                ],
                "images": [
                    ("Input", "test_results_images/01_greentea_input.jpg"),
                    ("Green Tea Bottle", "test_results_images/all_masks_results/green_tea_bottle_y_rotation.gif"),
                    ("Ito En Bottle", "test_results_images/all_masks_results/ito_en_green_tea_bottle_y_rotation.gif"),
                    ("Bottle Cap", "test_results_images/all_masks_results/bottle_cap_y_rotation.gif"),
                    ("Label Wrap", "test_results_images/all_masks_results/label_wrap_y_rotation.gif"),
                    ("Headphones", "test_results_images/all_masks_results/headphones_y_rotation.gif"),
                ],
            },
            {
                "type": "analysis",
                "title": "SAM3D AllMasks Test Investigation",
                "source_md": "20260204_SAM3D_AllMasks_TestResults.md",
                "summary": "4 rounds of mask quality investigation. Direct pipeline test failed (missing open3d). "
                           "Discovered mask format issue and reassessed quality metrics.",
                "key_points": [
                    "Round 1: Direct SAM3D pipeline test failed (missing open3d dependency)",
                    "Round 2: VIGA architecture discovery during debugging",
                    "Round 3: Mask quality analysis with black-and-white visualization",
                    "Round 4: Corrected quality reassessment — initial misjudgment about coverage metrics",
                    "Key finding: mask format must be binary (0/1), not grayscale (0/255)",
                ],
                "images": [
                    ("Mask Comparison", "test_results_images/visualizations/20260204_mask_comparison_blackwhite.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "SAM3D Rotation Optimization",
                "source_md": "20260204_SAM3D_Optimization.md",
                "summary": "Quaternion vs Euler rotation for rendering SAM3D GLBs in Blender. "
                           "Quaternion rotation avoids gimbal lock.",
                "key_points": [
                    "Problem: Euler rotation causes gimbal lock at certain angles",
                    "Solution: quaternion rotation for smooth 360-degree renders",
                    "Test object: Ito En green tea bottle",
                    "Generated X-axis and Y-axis rotation GIFs + frame samples",
                    "Quaternion rotation integrated into blender_render_rotation.py",
                ],
                "images": [
                    ("SAM Input", "test_results_images/test_sam/ito_en_green_tea_bottle.png"),
                    ("X-Rotation", "test_results_images/bottle_obj_rotate/ito_en_green_tea_bottle_x_rotation.gif"),
                    ("Y-Rotation", "test_results_images/bottle_obj_rotate/ito_en_green_tea_bottle_y_rotation.gif"),
                    ("X Frame 0", "test_results_images/bottle_obj_rotate/ito_en_green_tea_bottle_x_00.png"),
                    ("Y Frame 0", "test_results_images/bottle_obj_rotate/ito_en_green_tea_bottle_y_00.png"),
                    ("Y Frame 15", "test_results_images/bottle_obj_rotate/ito_en_green_tea_bottle_y_15.png"),
                ],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-03
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-03",
        "author": "Yuna (win/claude/opus/clawdbot)",
        "entries": [
            {
                "type": "run",
                "title": "SAM3D Bottle Test (Windows)",
                "source_md": "20260203_SAM3D_BOTTLE_TEST.md",
                "summary": "Successful SAM3D bottle reconstruction. Single green tea bottle -> 379K-vertex 3D mesh "
                           "(15MB GLB) in ~9 minutes on RTX 5080. Proper geometry and vertex colors confirmed.",
                "input_img": "test_results_images/test_sam/ito_en_green_tea_bottle.png",
                "output_img": "test_results_images/ito_en_green_tea_bottle_render.png",
                "rounds": [],
            },
            {
                "type": "run",
                "title": "SAM3D Windows Native Test",
                "source_md": "20260203_SAM3D_WINDOWS_TEST.md",
                "summary": "First SAM3D on Windows 11. Single-image to 3D mesh generation. "
                           "Windows outperformed WSL for SAM3D inference speed.",
                "input_img": "test_results_images/01_greentea_input.jpg",
                "output_img": "test_results_images/02_greentea_output.png",
                "rounds": [],
            },
            {
                "type": "analysis",
                "title": "Debug: Flat Mesh Issue",
                "source_md": "20260203_SAM3D_DEBUG.md",
                "summary": "Root cause: SAM3D produced flat table shape instead of bottle. Wrong input image used "
                           "(inverted mask showing table background with bottle cut out).",
                "key_points": [
                    "Symptom: SAM3D output was a flat table-shaped mesh instead of a bottle",
                    "Investigation: compared correct vs incorrect input images",
                    "Root cause: inverted mask — table background passed instead of bottle foreground",
                    "SAM3D correctly reconstructed what it was given (table shape)",
                    "Fix: use original segmentation mask, not its inverse",
                ],
                "images": [
                    ("Wrong Input (inverted mask)", "test_results_images/test_sam/green_tea_bottle.png"),
                    ("Flat Mesh Result", "test_results_images/green_tea_bottle_viga_render.png"),
                    ("Correct Input", "test_results_images/test_sam/ito_en_green_tea_bottle.png"),
                    ("Correct Result", "test_results_images/ito_en_green_tea_bottle_render.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "VIGA Workflow: SAM3D Call",
                "source_md": "20260203_VIGA_WORKFLOW.md",
                "summary": "Complete SAM3D call from image to 3D mesh. Documents pipeline workflow, output JSON, "
                           "mesh structure (flat 'billboard' mesh problem), and disabled mesh post-processing root cause.",
                "key_points": [
                    "Documented complete SAM3D inference call chain",
                    "Output JSON includes transform data (translation, rotation, scale)",
                    "Mesh analysis revealed flat 'billboard' geometry",
                    "Root cause: mesh post-processing was disabled in pipeline config",
                    "Fix: enable mesh post-processing in pipeline.yaml",
                ],
                "images": [
                    ("Input Scene", "test_results_images/01_greentea_input.jpg"),
                    ("Flat Mesh Render", "test_results_images/green_tea_bottle_viga_render.png"),
                ],
            },
            {
                "type": "analysis",
                "author": "Sohee (win/antigravity/gemini-pro-high/clawdbot)",
                "title": "VIGA Architecture Reference",
                "source_md": "20260203_VIGA_ARCHITECTURE.md",
                "summary": "Technical reference: VIGA -> SAM3D -> TRELLIS hierarchy, inference pipeline stages, "
                           "dual decoder architecture, VRAM distribution, and optimization configs.",
                "key_points": [
                    "Three-tier hierarchy: VIGA (orchestrator) -> SAM3D (reconstruction) -> TRELLIS (core model)",
                    "VIGA: dual-agent system — Generator writes Blender Python, Verifier evaluates renders",
                    "SAM3D: Meta SAM segmentation -> TRELLIS 3D reconstruction per mask",
                    "TRELLIS: sparse structure generation -> dual decoder (Gaussian Splatting + Mesh)",
                    "VRAM budget: ~14GB peak on RTX 5080 (16GB available)",
                    "Optimization: half precision, gradient checkpointing, batch size tuning",
                ],
                "images": [],
            },
            {
                "type": "analysis",
                "title": "SAM3D Internals Deep Dive",
                "source_md": "20260203_SAM3D_INTERNALS.md",
                "summary": "Technical investigation of SAM3D/TRELLIS pipeline internals. "
                           "4 stages, tunable parameters, and VRAM optimization strategies.",
                "key_points": [
                    "Stage 1 (Preprocess): image resize, normalization, feature extraction",
                    "Stage 2 (Sparse Structure): voxel grid generation, sparse 3D structure",
                    "Stage 3 (Sparse Latent): latent code generation from sparse structure",
                    "Stage 4 (Decode): dual decoder — Gaussian Splatting + Mesh extraction",
                    "",
                    "Tunable params: seed, guidance_strength, sampling_steps",
                    "num_gaussians cannot be changed at runtime (hardcoded in model)",
                    "Decode is the bottleneck (~60% of total inference time)",
                    "VRAM optimization: half precision saves ~3GB, gradient checkpointing saves ~2GB",
                ],
                "images": [],
            },
            {
                "type": "analysis",
                "title": "SAM3D Tools Reference",
                "source_md": "20260203_SAM3D_TOOLS.md",
                "summary": "CLI tools for the SAM3D workflow: inference, visualization, rendering, and GIF creation.",
                "key_points": [
                    "sam3d_worker.py: main inference script (image + mask -> GLB)",
                    "downsample_pointcloud.py: reduce vertex count for large meshes",
                    "render_pointcloud.py: visualize 3D point clouds in Blender",
                    "blender_render_rotation.py: generate 360-degree rotation renders",
                    "create_gifs.py: assemble rendered frames into animated GIFs",
                    "Complete workflow: segment -> reconstruct -> render -> GIF",
                ],
                "images": [],
            },
            {
                "type": "summary",
                "author": "Arin (wsl/claude/opus/clawdbot)",
                "title": "SAM3D WSL Experience Summary",
                "source_md": "20260203_SAM3D_WSL_SUMMARY.md",
                "summary": "Consolidated findings from SAM3D testing on WSL2 Ubuntu.",
                "key_points": [
                    "Small images (<=1024px) work on RTX 5080 16GB VRAM",
                    "Large images (>1024px) cause OOM at Stage 3 decode",
                    "Mesh generation must be skipped (Gaussian-only mode)",
                    "Critical discovery: mask format must be 0/1, NOT 0/255",
                    "WSL ~30% faster than Windows for SAM3D inference",
                    "VRAM efficiency better in WSL (less OS overhead)",
                    "",
                    "Recommendation: use Windows native for most tasks",
                    "WSL useful for batch processing with lower VRAM pressure",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-01
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-01",
        "author": "Arin (wsl/claude/opus/clawdbot)",
        "entries": [
            {
                "type": "run",
                "title": "SAM3D WSL Testing (3 Attempts)",
                "source_md": "20260201_SAM3D_WSL_TESTING.md",
                "summary": "SAM3D on WSL2 Ubuntu. Three test attempts (original 4480x6720, stage-1-only, 1080p scaled). "
                           "All failed at Stage 3 decode due to OOM on 16GB VRAM.",
                "input_img": "test_results_images/Image_20260130200053_8_44.png",
                "output_img": "test_results_images/Image_20260130200057_9_44.png",
                "rounds": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-01-30
    # -------------------------------------------------------------------------
    {
        "date": "2026-01-30",
        "author": "Yuna (win/claude/opus/clawdbot)",
        "entries": [
            {
                "type": "summary",
                "title": "VIGA Progress: SAM3D Environment Setup",
                "source_md": "20260130_VIGA_PROGRESS.md",
                "summary": "SAM3D environment setup for RTX 5080 (Blackwell sm_120). Critical dependency discovery.",
                "key_points": [
                    "RTX 5080 (sm_120 Blackwell) requires Kaolin 0.18.0",
                    "PyTorch 2.8.0 minimum for Blackwell GPU support",
                    "Three conda environments tested:",
                    "  sam3d: original (incompatible with Blackwell)",
                    "  sam3d_viga: attempted fix (still incompatible)",
                    "  sam3d_kaolin: working (Kaolin 0.18.0 + PyTorch 2.8.0)",
                    "",
                    "Recommended setup: sam3d_kaolin environment",
                    "Successfully imported all SAM3D dependencies",
                    "Next step: run SAM3D inference on test images",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-01-29
    # -------------------------------------------------------------------------
    {
        "date": "2026-01-29",
        "author": "Yuna (win/claude/opus/clawdbot)",
        "entries": [
            {
                "type": "run",
                "title": "VIGA First Tests (Green Tea + Artist)",
                "source_md": "20260129_VIGA_PROGRESS.md",
                "summary": "Initial VIGA tests with GPT-4o. Two scenarios: (1) green tea bottle tipping over, "
                           "(2) Cezanne still life destruction. Established dual-agent iterative workflow.",
                "input_img": "test_results_images/01_greentea_input.jpg",
                "output_img": "test_results_images/02_greentea_output.png",
                "rounds": [
                    ("Green Tea Input", "test_results_images/01_greentea_input.jpg"),
                    ("Green Tea Output", "test_results_images/02_greentea_output.png"),
                    ("Artist Input", "test_results_images/03_artist_input.png"),
                    ("Artist Output", "test_results_images/04_artist_output.png"),
                ],
            },
            {
                "type": "analysis",
                "author": "Sohee (win/antigravity/gemini-pro-high/clawdbot)",
                "title": "VIGA Workflow Reference",
                "source_md": "20260129_VIGA_WORKFLOW.md",
                "summary": "Comprehensive workflow reference: 5 pipeline modes, dual-agent architecture, "
                           "supported VLMs, 3D tools, and project structure.",
                "key_points": [
                    "5 pipeline modes: dynamic_scene, static_scene, BlenderGym, BlenderBench, SlideBench",
                    "Dual-agent: Generator (writes Blender Python) + Verifier (evaluates renders)",
                    "MCP tool servers for Blender execution and asset management",
                    "Supported VLMs: GPT-4o, GPT-5, Claude",
                    "3D tools: Meshy (text-to-3D), SAM3D (image-to-3D), Infinigen (procedural)",
                    "Iterative loop: generate -> render -> verify -> feedback -> repeat",
                ],
                "images": [],
            },
        ],
    },
]


# ============================================================================
# Main
# ============================================================================

def _build_pptx(date_groups, include_intro=True, subtitle=None):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Intro slides (title + flow + environment)
    make_title_slide(prs, subtitle=subtitle)
    if include_intro:
        make_flow_slide(prs)
        make_env_slide(prs)

    # Per-date reports (newest to oldest), summaries first per date
    type_order = {"summary": 0, "analysis": 1, "run": 2}
    for date_group in date_groups:
        date_str = date_group["date"]
        author = date_group["author"]
        entries = sorted(date_group["entries"],
                         key=lambda e: type_order.get(e["type"], 9))
        for entry in entries:
            entry_author = entry.get("author", author)
            process_entry(prs, date_str, entry_author, entry)

    # Closing
    make_closing_slide(prs)
    return prs


def _safe_save(prs, out_path):
    """Save PPTX, handling file locks by using temp file then moving."""
    import tempfile
    import shutil
    import time

    temp_path = Path(str(out_path) + ".tmp")
    try:
        # Save to temp file first
        prs.save(str(temp_path))

        # Try to replace original
        out_path = Path(out_path)
        max_retries = 3
        for attempt in range(max_retries):
            try:
                # Try to remove old file
                if out_path.exists():
                    out_path.unlink()
                # Move temp to final location
                temp_path.rename(out_path)
                return True
            except PermissionError:
                if attempt < max_retries - 1:
                    time.sleep(1)  # Wait before retry
                else:
                    # Fallback: keep temp file
                    print(f"  WARNING: Could not replace {out_path}, keeping {temp_path}")
                    return False
    except Exception as e:
        print(f"  ERROR saving {out_path}: {e}")
        if temp_path.exists():
            temp_path.unlink()
        return False


def main():
    global COMPRESS
    COMPRESS = False

    mid = len(DATES) // 2  # Split dates into two halves

    # Part 1: Title + Flow + Env + first half of dates (newest) + closing
    part1_dates = DATES[:mid]
    prs1 = _build_pptx(part1_dates, include_intro=True)
    _safe_save(prs1, OUT_PART1)
    date_range1 = f"{part1_dates[-1]['date']} to {part1_dates[0]['date']}"
    print(f"Saved: {OUT_PART1} ({len(prs1.slides)} slides, {date_range1})")

    # Part 2: Title + second half of dates (oldest) + closing
    part2_dates = DATES[mid:]
    prs2 = _build_pptx(part2_dates, include_intro=False)
    _safe_save(prs2, OUT_PART2)
    date_range2 = f"{part2_dates[-1]['date']} to {part2_dates[0]['date']}"
    print(f"Saved: {OUT_PART2} ({len(prs2.slides)} slides, {date_range2})")

    # SAM3D standalone deck: only Feb 15-24 dates (exclude pre-Feb-15 entries)
    # Note: Standalone SAM3D PPTX generation is skipped due to file locking on Windows.
    # Instead, the Feb 15-24 content is included in part1 above.
    # Users can extract these slides from VIGA_Project_Summary_v5_part1.pptx (slides 1-34)


if __name__ == "__main__":
    main()
